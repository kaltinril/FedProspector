"""Generate FedProspect ETL Architecture PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MEDIUM_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
BLACK = RGBColor(0x00, 0x00, 0x00)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
ACCENT_GREEN = RGBColor(0x38, 0x8E, 0x3C)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_bar(slide, title_text):
    """Add a dark blue title bar at the top of a slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, bullet=False):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if isinstance(text, list):
        for i, line in enumerate(text):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = alignment
            if bullet:
                p.level = 0
                p.space_before = Pt(4)
    else:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
    return txBox


def add_table(slide, left, top, width, height, rows, col_widths=None):
    """Add a table with header row styled dark blue."""
    shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = BLACK

            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

            cell.text_frame.margin_left = Inches(0.05)
            cell.text_frame.margin_right = Inches(0.05)
            cell.text_frame.margin_top = Inches(0.03)
            cell.text_frame.margin_bottom = Inches(0.03)

    return shape


def add_flow_box(slide, left, top, width, height, text, fill_color=ACCENT_BLUE,
                 font_size=12, font_color=WHITE):
    """Add a rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top):
    """Add a connector arrow between two points."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = RGBColor(0x66, 0x66, 0x66)
    connector.line.width = Pt(2)
    return connector


# =====================================================================
# SLIDE 1: System Overview
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_title_bar(slide, "FedProspect ETL Architecture - System Overview")

# Data sources column
add_text_box(slide, Inches(0.3), Inches(1.3), Inches(2), Inches(0.4),
             "Data Sources", font_size=16, bold=True, color=DARK_BLUE)
sources = [
    ("SAM.gov\nOpportunities", ACCENT_BLUE),
    ("SAM.gov\nAwards (FPDS)", ACCENT_BLUE),
    ("SAM.gov\nEntities", ACCENT_BLUE),
    ("USASpending.gov\nBulk Awards", ACCENT_BLUE),
    ("GSA CALC+\nLabor Rates", ACCENT_BLUE),
]
for i, (label, color) in enumerate(sources):
    add_flow_box(slide, Inches(0.3), Inches(1.8 + i * 1.05), Inches(2.2), Inches(0.85),
                 label, fill_color=color, font_size=11)

# ETL column
add_text_box(slide, Inches(3.5), Inches(1.3), Inches(3), Inches(0.4),
             "Python ETL Pipeline", font_size=16, bold=True, color=DARK_BLUE)
etl_steps = [
    ("API Clients\nRate-limited, retries", RGBColor(0x5B, 0x9B, 0xD5)),
    ("Staging Tables\nRaw JSON + SHA-256", RGBColor(0x5B, 0x9B, 0xD5)),
    ("Change Detection\nInsert / Update / Skip", RGBColor(0x5B, 0x9B, 0xD5)),
    ("Data Cleaning\n10 quality rules", RGBColor(0x5B, 0x9B, 0xD5)),
    ("Load Manager\netl_load_log tracking", RGBColor(0x5B, 0x9B, 0xD5)),
]
for i, (label, color) in enumerate(etl_steps):
    add_flow_box(slide, Inches(3.5), Inches(1.8 + i * 1.05), Inches(2.8), Inches(0.85),
                 label, fill_color=color, font_size=11)

# Database column
add_text_box(slide, Inches(7.3), Inches(1.3), Inches(2.5), Inches(0.4),
             "MySQL 8.4 Database", font_size=16, bold=True, color=DARK_BLUE)
db_items = [
    ("Opportunities\n+ Resource Links", ACCENT_GREEN),
    ("Awards\n+ Transactions", ACCENT_GREEN),
    ("Entities\n+ Certifications", ACCENT_GREEN),
    ("Attachments\n+ Extracted Intel", ACCENT_GREEN),
    ("Reference Data\nNAICS, PSC, FIPS", ACCENT_GREEN),
]
for i, (label, color) in enumerate(db_items):
    add_flow_box(slide, Inches(7.3), Inches(1.8 + i * 1.05), Inches(2.5), Inches(0.85),
                 label, fill_color=color, font_size=11)

# App layer column
add_text_box(slide, Inches(10.6), Inches(1.3), Inches(2.5), Inches(0.4),
             "Application Layer", font_size=16, bold=True, color=DARK_BLUE)
app_items = [
    ("C# ASP.NET Core\nREST API", ACCENT_ORANGE),
    ("React 19 + TS\nVite 8 + MUI v7", ACCENT_ORANGE),
    ("Prospect Mgmt\npWin, Qualification", ACCENT_ORANGE),
]
for i, (label, color) in enumerate(app_items):
    add_flow_box(slide, Inches(10.6), Inches(1.8 + i * 1.4), Inches(2.5), Inches(1.1),
                 label, fill_color=color, font_size=11)

# Arrows between columns
for i in range(5):
    y = Inches(2.2 + i * 1.05)
    add_arrow(slide, Inches(2.55), y, Inches(3.45), y)
    add_arrow(slide, Inches(6.35), y, Inches(7.25), y)

for i in range(3):
    y = Inches(2.35 + i * 1.4)
    add_arrow(slide, Inches(9.85), y, Inches(10.55), y)

# =====================================================================
# SLIDE 2: Data Source Inventory
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Data Source Inventory")

table_data = [
    ["Source", "API Type", "Rate Limit", "Loader Module", "Frequency", "Key Data"],
    ["SAM.gov Opportunities", "REST (GET)", "Key1: 10/day\nKey2: 1000/day",
     "opportunity_loader.py", "Daily (31d back)", "Solicitations, notices"],
    ["SAM.gov Awards (FPDS)", "REST (GET)", "Key1: 10/day\nKey2: 1000/day",
     "awards_loader.py", "Daily (10d back)", "Contract awards by set-aside"],
    ["SAM.gov Entities", "Bulk DAT file", "N/A (file download)", "entity_loader.py\nbulk_loader.py",
     "Monthly full /\nDaily incremental", "Vendor registrations"],
    ["SAM.gov Exclusions", "REST (GET)", "Key1: 10/day\nKey2: 1000/day",
     "exclusion_loader.py", "Weekly", "Debarred/suspended vendors"],
    ["USASpending.gov", "REST (POST)", "No hard limit", "usaspending_bulk_loader.py",
     "Daily (5d back)", "Award amounts, transactions"],
    ["USASpending Subawards", "REST (POST)", "No hard limit", "subaward_loader.py",
     "On-demand", "Subcontract data"],
    ["GSA CALC+", "REST (GET)", "No hard limit", "calc_loader.py", "On-demand",
     "Labor rates by category"],
    ["Federal Hierarchy", "REST (GET)", "Key2: 1000/day", "fedhier_loader.py",
     "On-demand", "Agency/office hierarchy"],
]

add_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.5), table_data,
          col_widths=[Inches(2.2), Inches(1.3), Inches(1.8), Inches(2.2), Inches(1.8), Inches(2.2)])

# =====================================================================
# SLIDE 3: Daily Load Pipeline
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Daily Load Pipeline (daily_load.bat)")

steps = [
    ("1", "Load Opportunities", "31 days back, key 2, max 300 calls, --force"),
    ("2", "Fetch Descriptions", "Priority NAICS+set-aside first, limit 100/day"),
    ("3", "Load USASpending Bulk", "5 days back, award + transaction data"),
    ("4", "Load Awards - 8(a)", "10 days back, 24 NAICS codes, key 2"),
    ("5", "Load Awards - WOSB", "10 days back, 24 NAICS codes, key 2"),
    ("6", "Load Awards - SBA", "10 days back, 24 NAICS codes, key 2"),
    ("7", "Enrich Link Metadata", "Filenames, content types for resource links"),
    ("8", "Download Attachments", "Active opps, missing only, batch 5000"),
    ("9", "Extract Text", "PDF/DOCX/XLSX text extraction, batch 5000"),
    ("10", "Extract Intel", "Keyword intelligence from text, batch 5000"),
    ("11", "Backfill Opp Intel", "Propagate intel findings to opportunity table"),
    ("12", "Cleanup Files", "Delete files after all stages complete"),
]

# Left column (steps 1-6)
for i, (num, title, desc) in enumerate(steps[:6]):
    y = Inches(1.4 + i * 0.95)
    add_flow_box(slide, Inches(0.3), y, Inches(0.6), Inches(0.55),
                 num, fill_color=DARK_BLUE, font_size=16)
    add_flow_box(slide, Inches(1.0), y, Inches(2.0), Inches(0.55),
                 title, fill_color=ACCENT_BLUE, font_size=11)
    add_text_box(slide, Inches(3.1), y + Inches(0.05), Inches(3.2), Inches(0.5),
                 desc, font_size=10, color=BLACK)

# Right column (steps 7-12)
for i, (num, title, desc) in enumerate(steps[6:]):
    y = Inches(1.4 + i * 0.95)
    add_flow_box(slide, Inches(6.8), y, Inches(0.6), Inches(0.55),
                 num, fill_color=DARK_BLUE, font_size=16)
    add_flow_box(slide, Inches(7.5), y, Inches(2.3), Inches(0.55),
                 title, fill_color=ACCENT_BLUE, font_size=11)
    add_text_box(slide, Inches(9.9), y + Inches(0.05), Inches(3.2), Inches(0.5),
                 desc, font_size=10, color=BLACK)

add_text_box(slide, Inches(0.3), Inches(6.9), Inches(12), Inches(0.4),
             "NAICS filter: 336611, 488190, 519210, 541219, 541330, 541511-541990, 561110-561990, 611430-624190, 812910",
             font_size=10, color=RGBColor(0x66, 0x66, 0x66))

# =====================================================================
# SLIDE 4: Attachment Intelligence Pipeline
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Attachment Intelligence Pipeline (7-Stage State Machine)")

stages = [
    ("1. Load", "Opportunities loaded\nwith resource_link\nURLs from SAM.gov",
     RGBColor(0x2E, 0x75, 0xB6)),
    ("2. Enrich", "Fetch filename &\ncontent type via\nHEAD requests",
     RGBColor(0x2E, 0x86, 0xC1)),
    ("3. Download", "Download files to\nE:\\fedprospector\\\nattachments\\",
     RGBColor(0x27, 0xAE, 0x60)),
    ("4. Extract", "PDF, DOCX, XLSX\ntext extraction\n(parallel processing)",
     RGBColor(0x27, 0xAE, 0x60)),
    ("5. Keywords", "Keyword intelligence\nextraction from\nattachment text",
     ACCENT_ORANGE),
    ("6. AI Analysis", "LLM-based analysis\nof extracted text\n(Phase 110C)",
     ACCENT_ORANGE),
    ("7. Cleanup", "Delete source files\nafter ALL stages\ncomplete",
     RGBColor(0xC0, 0x39, 0x2B)),
]

box_w = Inches(1.6)
gap = Inches(0.15)
start_x = Inches(0.3)

for i, (title, desc, color) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    add_flow_box(slide, x, Inches(1.5), box_w, Inches(0.5),
                 title, fill_color=color, font_size=11)
    add_text_box(slide, x, Inches(2.1), box_w, Inches(1.2),
                 desc, font_size=10, color=BLACK, alignment=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_arrow(slide, x + box_w, Inches(1.75), x + box_w + gap, Inches(1.75))

# Key details
details = [
    "Storage: E:\\fedprospector\\attachments\\ (configurable via ATTACHMENT_DIR env var)",
    "Supported formats: PDF (PyPDF2 + find_tables), DOCX, XLSX (openpyxl), DOC (LibreOffice)",
    "Parallel text extraction for multi-page PDFs (>30 pages skip table detection)",
    "Cleanup is gated: files only deleted after ALL stages complete for that attachment",
    "Magic byte detection validates file type regardless of extension",
    "Batch processing with configurable batch_size (default 5000)",
]
for i, detail in enumerate(details):
    add_text_box(slide, Inches(0.5), Inches(3.5 + i * 0.45), Inches(12), Inches(0.4),
                 f"  {detail}", font_size=12, color=BLACK)

# =====================================================================
# SLIDE 5: Change Detection
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Change Detection - SHA-256 Record Hashing")

# Flow diagram
flow_steps = [
    ("API Response\n(raw JSON)", RGBColor(0x2E, 0x75, 0xB6)),
    ("Stage to\nstg_*_raw table", RGBColor(0x5B, 0x9B, 0xD5)),
    ("SHA-256 hash\nof key fields", RGBColor(0x5B, 0x9B, 0xD5)),
    ("Compare to\nexisting hash", RGBColor(0xE6, 0x7E, 0x22)),
]
for i, (label, color) in enumerate(flow_steps):
    x = Inches(0.5 + i * 3.1)
    add_flow_box(slide, x, Inches(1.5), Inches(2.5), Inches(0.8),
                 label, fill_color=color, font_size=13)
    if i < len(flow_steps) - 1:
        add_arrow(slide, x + Inches(2.5), Inches(1.9), x + Inches(3.1), Inches(1.9))

# Three outcomes
outcomes = [
    ("NEW RECORD\n(not in hash cache)", "INSERT into\ntarget table", ACCENT_GREEN),
    ("HASH DIFFERS\n(record changed)", "Compute field diff\nLog to *_history\nUPDATE record", ACCENT_ORANGE),
    ("HASH MATCHES\n(unchanged)", "SKIP\n(no DB write)", RGBColor(0x95, 0xA5, 0xA6)),
]
for i, (condition, action, color) in enumerate(outcomes):
    x = Inches(0.5 + i * 4.2)
    add_flow_box(slide, x, Inches(2.8), Inches(3.5), Inches(0.7),
                 condition, fill_color=color, font_size=12)
    add_text_box(slide, x, Inches(3.6), Inches(3.5), Inches(0.8),
                 action, font_size=12, color=BLACK, alignment=PP_ALIGN.CENTER)

# Key details
details = [
    "Hash cache loaded once at load start for performance (single SELECT query)",
    "Staging tables (stg_*_raw): raw_json + raw_record_hash columns",
    "StagingMixin: shared base for 7 loaders (opportunities, awards, exclusions, etc.)",
    "Field-level diffs logged to *_history tables for audit trail",
    "UPSERT (INSERT ON DUPLICATE KEY UPDATE) ensures idempotent loads",
    "Hash fields are configurable per loader (e.g., 32 fields for entities)",
]
for i, detail in enumerate(details):
    add_text_box(slide, Inches(0.5), Inches(4.7 + i * 0.4), Inches(12), Inches(0.35),
                 f"  {detail}", font_size=12, color=BLACK)

# =====================================================================
# SLIDE 6: Rate Limiting & API Keys
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Rate Limiting & API Key Management")

# Key comparison table
key_data = [
    ["", "Key 1 (SAM_API_KEY)", "Key 2 (SAM_API_KEY_2)"],
    ["Daily Limit", "10 requests/day", "1,000 requests/day"],
    ["Tier", "Free tier", "Paid tier"],
    ["Tracked As", "SAM_KEY1", "SAM_KEY2"],
    ["Used By", "Manual / ad-hoc loads", "Daily pipeline (default)"],
    ["Shared Pool", "Yes - all SAM endpoints", "Yes - all SAM endpoints"],
]
add_table(slide, Inches(0.5), Inches(1.4), Inches(7), Inches(3.0), key_data,
          col_widths=[Inches(1.8), Inches(2.6), Inches(2.6)])

# Rate limit mechanism
add_text_box(slide, Inches(8), Inches(1.4), Inches(5), Inches(0.4),
             "Rate Limit Tracking", font_size=16, bold=True, color=DARK_BLUE)

mechanisms = [
    "etl_rate_limit table: source_system + request_date",
    "INSERT ON DUPLICATE KEY UPDATE on each API call",
    "Pre-request check: _get_remaining_requests()",
    "429 response: parse nextAccessTime, raise immediately",
    "Connection pooling: 10 connections, reuse TCP",
    "Inter-request delay: 0.1s (configurable)",
    "Exponential backoff on 5xx: 2^attempt seconds",
    "Max 3 retries per request (configurable)",
    "Connection/timeout errors also trigger retries",
]
for i, item in enumerate(mechanisms):
    add_text_box(slide, Inches(8.2), Inches(1.9 + i * 0.42), Inches(4.8), Inches(0.4),
                 f"  {item}", font_size=12, color=BLACK)

# Pagination
add_text_box(slide, Inches(0.5), Inches(4.7), Inches(5), Inches(0.4),
             "Generic Pagination Engine", font_size=16, bold=True, color=DARK_BLUE)
pag_items = [
    "3 strategies: offset, page-number, has_next_key",
    "Safety guard: max_pages=1000 (prevents infinite loops)",
    "Generator-based: yields pages lazily",
    "Supports: total_key, total_pages_key, has_next_key",
    "USASpending: page-based with has_next",
    "SAM.gov: offset-based with totalRecords",
]
for i, item in enumerate(pag_items):
    add_text_box(slide, Inches(0.7), Inches(5.2 + i * 0.35), Inches(5), Inches(0.3),
                 f"  {item}", font_size=11, color=BLACK)

# =====================================================================
# SLIDE 7: Error Handling
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Error Handling & Resumable Loads")

# Load log tracking
add_text_box(slide, Inches(0.3), Inches(1.3), Inches(4), Inches(0.4),
             "etl_load_log Table", font_size=18, bold=True, color=DARK_BLUE)

log_data = [
    ["Column", "Purpose"],
    ["load_id", "Auto-increment primary key"],
    ["source_system", "SAM_OPPORTUNITY, SAM_AWARD, etc."],
    ["load_type", "FULL, INCREMENTAL, DAILY"],
    ["status", "RUNNING, SUCCESS, FAILED"],
    ["started_at / completed_at", "Load duration tracking"],
    ["records_read/inserted/updated", "Per-load counters"],
    ["records_errored", "Failed record count"],
    ["parameters (JSON)", "Date range, pages_fetched, complete flag"],
    ["error_message", "Failure reason (truncated to 5000 chars)"],
]
add_table(slide, Inches(0.3), Inches(1.8), Inches(5.5), Inches(4.5), log_data,
          col_widths=[Inches(2.5), Inches(3.0)])

# Resumable loads
add_text_box(slide, Inches(6.5), Inches(1.3), Inches(6), Inches(0.4),
             "Resumable Load Flow", font_size=18, bold=True, color=DARK_BLUE)

resume_steps = [
    "1. save_load_progress() called after each API page",
    "2. Parameters JSON stores pages_fetched, complete=false",
    "3. If process crashes, load stays status=SUCCESS",
    "4. Next run: get_resumable_load() finds incomplete load",
    "5. Resumes from last checkpointed page",
    "6. On completion: parameters.complete = true",
]
for i, step in enumerate(resume_steps):
    add_text_box(slide, Inches(6.7), Inches(1.9 + i * 0.45), Inches(6), Inches(0.4),
                 step, font_size=13, color=BLACK)

# Stale cleanup
add_text_box(slide, Inches(6.5), Inches(4.8), Inches(6), Inches(0.4),
             "Additional Safeguards", font_size=16, bold=True, color=DARK_BLUE)

safeguards = [
    "cleanup_stale_running(): auto-fail loads stuck >2 hours",
    "etl_load_error: per-record errors with raw_data (10KB limit)",
    "Staging rows marked 'E' with error_message on failure",
    "Staging rows marked 'Y' on successful normalization",
    "Over-counting rate limits preferred over under-counting",
]
for i, item in enumerate(safeguards):
    add_text_box(slide, Inches(6.7), Inches(5.3 + i * 0.4), Inches(6), Inches(0.35),
                 f"  {item}", font_size=12, color=BLACK)

# =====================================================================
# SLIDE 8: Data Quality Rules
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Data Quality Rules Engine")

# Hardcoded rules
add_text_box(slide, Inches(0.3), Inches(1.3), Inches(6), Inches(0.4),
             "10 Built-in Quality Rules (DataCleaner)", font_size=16, bold=True, color=DARK_BLUE)

rules_data = [
    ["#", "Rule", "Action"],
    ["1", "ZIP codes with city/state/country names", "Extract 5-digit or 5+4 ZIP"],
    ["2", "ZIP codes with PO BOX data", "Extract ZIP, discard PO BOX"],
    ["3", "State fields containing dates", "Clear state, preserve date"],
    ["4", "Foreign provinces in state field", "Flag only (no auto-fix)"],
    ["5", "Non-ASCII chars in country names", "Unicode normalization"],
    ["6", "Missing country codes (XKS/XWB/XGZ)", "Validate only"],
    ["7", "Comma-separated CAGE codes", "Split and normalize"],
    ["8", "Retired NAICS codes", "Flag only"],
    ["9", "Escaped pipes in DAT files", "Unescape pipe characters"],
    ["10", "YYYYMMDD date strings", "Convert to DATE type"],
]
add_table(slide, Inches(0.3), Inches(1.8), Inches(6.3), Inches(5.0), rules_data,
          col_widths=[Inches(0.4), Inches(3.0), Inches(2.9)])

# DB rules
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(6), Inches(0.4),
             "etl_data_quality_rule Table", font_size=16, bold=True, color=DARK_BLUE)

db_features = [
    "Configurable rules stored in database (not hardcoded)",
    "Rules loaded at DataCleaner initialization",
    "Each rule has: rule_name, table_name, column_name",
    "rule_type: REGEX, LOOKUP, RANGE, CUSTOM",
    "rule_expression: regex pattern or validation logic",
    "action: CLEAN, FLAG, REJECT, LOG",
    "severity: INFO, WARNING, ERROR",
    "is_active flag for enable/disable without code changes",
    "Stats tracking: counts per rule application",
    "Extensible: add new rules via INSERT, no deploy needed",
]
for i, item in enumerate(db_features):
    add_text_box(slide, Inches(7.2), Inches(1.9 + i * 0.42), Inches(5.8), Inches(0.4),
                 f"  {item}", font_size=12, color=BLACK)

# =====================================================================
# SLIDE 9: Bulk Loading (LOAD DATA INFILE)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "Bulk Loading - LOAD DATA INFILE")

# Flow
flow = [
    ("SAM.gov\nMonthly\nDAT File", RGBColor(0x2E, 0x75, 0xB6)),
    ("DAT Parser\nPipe-delimited\nField extraction", RGBColor(0x5B, 0x9B, 0xD5)),
    ("Data Cleaning\n10 quality rules\nNormalization", RGBColor(0x5B, 0x9B, 0xD5)),
    ("TSV Temp File\nEscaped values\nNULL handling", ACCENT_ORANGE),
    ("LOAD DATA\nINFILE\nMySQL bulk load", ACCENT_GREEN),
]
for i, (label, color) in enumerate(flow):
    x = Inches(0.3 + i * 2.55)
    add_flow_box(slide, x, Inches(1.5), Inches(2.2), Inches(1.0),
                 label, fill_color=color, font_size=12)
    if i < len(flow) - 1:
        add_arrow(slide, x + Inches(2.2), Inches(2.0), x + Inches(2.55), Inches(2.0))

# Details
add_text_box(slide, Inches(0.3), Inches(2.8), Inches(6), Inches(0.4),
             "Why LOAD DATA INFILE?", font_size=16, bold=True, color=DARK_BLUE)

why_items = [
    "10-100x faster than INSERT loops for large datasets",
    "SAM.gov entity file: ~500K+ records per monthly extract",
    "Single LOAD DATA statement replaces thousands of INSERTs",
    "MySQL secure-file-priv=\"\" allows loading from any path",
    "fed_app user has FILE privilege for LOAD DATA access",
]
for i, item in enumerate(why_items):
    add_text_box(slide, Inches(0.5), Inches(3.3 + i * 0.38), Inches(5.5), Inches(0.35),
                 f"  {item}", font_size=12, color=BLACK)

add_text_box(slide, Inches(7.0), Inches(2.8), Inches(6), Inches(0.4),
             "Implementation Details", font_size=16, bold=True, color=DARK_BLUE)

impl_items = [
    "bulk_loader.py: Entity-specific LOAD DATA pipeline",
    "Writes temp TSV files with escape_tsv_value()",
    "32 entity hash fields for change detection",
    "Column order matches entity table exactly",
    "Uses tempfile module for auto-cleanup",
    "LoadManager tracks records_read/inserted/updated",
    "Also used by: calc_loader, usaspending_bulk_loader",
    "USASpending: parallel downloads + bulk DB loading",
]
for i, item in enumerate(impl_items):
    add_text_box(slide, Inches(7.2), Inches(3.3 + i * 0.38), Inches(5.8), Inches(0.35),
                 f"  {item}", font_size=12, color=BLACK)

# =====================================================================
# SLIDE 10: On-Demand Loading
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "On-Demand Loading (data_load_request)")

# Flow
flow_steps_demand = [
    ("React UI\nUser action\ntriggers request", RGBColor(0xE6, 0x7E, 0x22)),
    ("C# API\nInserts row\nstatus=PENDING", RGBColor(0xE6, 0x7E, 0x22)),
    ("data_load_request\nTable\n(queue)", RGBColor(0x2E, 0x75, 0xB6)),
    ("DemandLoader\nPolls for\nPENDING rows", RGBColor(0x5B, 0x9B, 0xD5)),
    ("API Client\nFetches from\nexternal source", RGBColor(0x5B, 0x9B, 0xD5)),
    ("ETL Loader\nNormalizes &\nstores data", ACCENT_GREEN),
]
for i, (label, color) in enumerate(flow_steps_demand):
    x = Inches(0.2 + i * 2.15)
    add_flow_box(slide, x, Inches(1.5), Inches(1.9), Inches(1.0),
                 label, fill_color=color, font_size=11)
    if i < len(flow_steps_demand) - 1:
        add_arrow(slide, x + Inches(1.9), Inches(2.0), x + Inches(2.15), Inches(2.0))

# Request types
add_text_box(slide, Inches(0.3), Inches(2.9), Inches(6), Inches(0.4),
             "Supported Request Types", font_size=16, bold=True, color=DARK_BLUE)

type_data = [
    ["Request Type", "Source", "Loader Used", "Description"],
    ["USASPENDING_AWARD", "USASpending.gov", "USASpendingLoader", "Fetch award + transactions by ID"],
    ["FPDS_AWARD", "SAM.gov Awards", "AwardsLoader", "Fetch FPDS contract data by ID"],
    ["REFRESH_FEDHIER_ORG", "SAM.gov Fed Hierarchy", "FedHierLoader", "Refresh single agency/office org"],
    ["ATTACHMENT_ANALYSIS", "Local processing", "AttachmentPipeline", "AI analysis of opp attachments"],
]
add_table(slide, Inches(0.3), Inches(3.4), Inches(12), Inches(2.2), type_data,
          col_widths=[Inches(2.5), Inches(2.0), Inches(2.5), Inches(5.0)])

# Status flow
add_text_box(slide, Inches(0.3), Inches(5.2), Inches(6), Inches(0.4),
             "Request Status Flow", font_size=16, bold=True, color=DARK_BLUE)

statuses = [
    ("PENDING", RGBColor(0x2E, 0x75, 0xB6)),
    ("PROCESSING", ACCENT_ORANGE),
    ("COMPLETED", ACCENT_GREEN),
    ("FAILED", RGBColor(0xC0, 0x39, 0x2B)),
]
for i, (status, color) in enumerate(statuses):
    x = Inches(0.5 + i * 2.5)
    add_flow_box(slide, x, Inches(5.7), Inches(1.8), Inches(0.6),
                 status, fill_color=color, font_size=14)
    if i < len(statuses) - 1:
        add_arrow(slide, x + Inches(1.8), Inches(6.0), x + Inches(2.5), Inches(6.0))

details = [
    "DemandLoader processes up to 10 pending requests per poll cycle",
    "Uses existing ETL loaders and API clients (no duplicate logic)",
    "Failed requests store error message for debugging",
    "SAM APIs use Key 2 (1000/day); Fed Hierarchy refresh = 1 API call per org",
]
for i, item in enumerate(details):
    add_text_box(slide, Inches(0.5), Inches(6.5 + i * 0.3), Inches(12), Inches(0.3),
                 f"  {item}", font_size=11, color=BLACK)

# Save
output_path = r"c:\git\fedProspect\docs\etl\FedProspect-ETL-Architecture.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
