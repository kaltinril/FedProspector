"""Generate FedProspect-Data-Flows.pptx documenting ETL pipeline data flows."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# -- Constants ---------------------------------------------------------------
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x00, 0x6E, 0xB8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_TOP = Inches(0.0)
TITLE_HEIGHT = Inches(1.0)
BODY_TOP = Inches(1.1)
BODY_LEFT = Inches(0.5)
BODY_WIDTH = Inches(12.333)
BODY_HEIGHT = Inches(6.0)


def add_title_bar(slide, title_text):
    """Add a dark-blue title bar across the top of the slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), TITLE_TOP, SLIDE_W, TITLE_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.2)


def add_body_text(slide, lines, *, font_size=14, line_spacing=1.4):
    """Add body text lines to the slide below the title bar."""
    txBox = slide.shapes.add_textbox(BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Determine styling based on prefix
        text = line
        is_heading = False
        is_note = False
        is_arrow = False
        indent_level = 0

        if line.startswith("## "):
            text = line[3:]
            is_heading = True
        elif line.startswith(">> "):
            text = line[3:]
            is_note = True
        elif line.startswith("   "):
            indent_level = 1
            text = line.lstrip()

        p.text = text

        if is_heading:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
            p.space_before = Pt(14)
        elif is_note:
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = MED_GRAY
            p.space_before = Pt(4)
        else:
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

        if indent_level > 0:
            p.level = indent_level
            p.font.size = Pt(font_size - 1)

        p.line_spacing = Pt(font_size * line_spacing)


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # ====================================================================
    # Slide 1: SAM.gov Opportunity Flow
    # ====================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "SAM.gov Opportunity Pipeline")
    add_body_text(slide, [
        "## Data Flow",
        "",
        "SAM.gov Opportunities API  -->  JSON response  -->  Parse & flatten nested fields",
        "   --> stg_opportunity_raw staging table  -->  SHA-256 hash compare (ChangeDetector)",
        "   --> UPSERT into opportunity table  -->  Extract resource_links JSON array",
        "",
        "## Key Details",
        "",
        "Source:           SAM.gov Opportunities API (v1/api/search, paged)",
        "Staging table:    stg_opportunity_raw (keyed on notice_id)",
        "Target table:     opportunity (PK: notice_id)",
        "Change detection: SHA-256 hash of 25 business fields, skip unchanged records",
        "Batch size:       500 records per upsert batch",
        "CLI command:      python main.py load opportunities [--key=1|2]",
        "",
        "## Volume",
        "",
        "~38,000 total opportunities loaded",
        "~13,000 opportunities have resource_links (attachments)",
        "",
        ">> Uses StagingMixin for load tracking. Filters by NAICS codes and set-aside types.",
    ])

    # ====================================================================
    # Slide 2: SAM.gov Award Flow (FPDS)
    # ====================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "SAM.gov Awards Pipeline (FPDS Contracts)")
    add_body_text(slide, [
        "## Data Flow",
        "",
        "SAM.gov Awards API  -->  JSON response  -->  Parse contract fields",
        "   --> stg_award staging  -->  SHA-256 hash compare (composite key: contract_id|mod_number)",
        "   --> UPSERT into fpds_contract table",
        "",
        "## Key Details",
        "",
        "Source:           SAM.gov Contract Awards API (requires API key 2, 1000/day)",
        "Target table:     fpds_contract (composite PK: contract_id, modification_number)",
        "Change detection: SHA-256 hash of 16 business fields",
        "Hash key:         contract_id|modification_number (pipe-separated composite)",
        "CLI command:      python main.py load awards [--key=2]",
        "",
        "## Filters Applied",
        "",
        "NAICS code filtering (target industry codes)",
        "Set-aside type filtering (WOSB, 8(a), etc.)",
        "Tracks: dollars_obligated, base_and_all_options, extent_competed, number_of_offers",
        "",
        ">> Uses StagingMixin + batch_upsert pattern. Composite PK requires special hash key handling.",
    ])

    # ====================================================================
    # Slide 3: SAM.gov Entity Flow
    # ====================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "SAM.gov Entity Pipeline (Vendor Registration)")
    add_body_text(slide, [
        "## Data Flow",
        "",
        "SAM.gov Bulk Extract  -->  JSON file (monthly/daily)  -->  ijson streaming parser",
        "   --> stg_entity staging (keyed on uei_sam, tracked by load_id)",
        "   --> SHA-256 hash compare  -->  UPSERT entity (parent)",
        "   --> DELETE + re-INSERT 8 child tables",
        "",
        "## Key Details",
        "",
        "Source:           SAM.gov Entity Bulk Extract (JSON, not API)",
        "Streaming:        ijson streaming parser for large files (avoids full JSON in memory)",
        "Target tables:    entity (PK: uei_sam) + 8 child tables",
        "Child tables:     entity_naics, entity_psc, entity_poc, entity_address,",
        "                  entity_certification, entity_goods_and_services,",
        "                  entity_sba_business_type, entity_url",
        "Change detection: SHA-256 hash of 31 business fields",
        "Batch size:       1,000 records per batch",
        "CLI command:      python main.py load entities",
        "",
        ">> Does NOT use StagingMixin -- custom streaming staging logic for memory efficiency.",
        ">> Field-level history tracking for audit trail on changed fields.",
    ])

    # ====================================================================
    # Slide 4: USASpending Bulk Flow
    # ====================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "USASpending.gov Bulk Pipeline")
    add_body_text(slide, [
        "## Data Flow",
        "",
        "USASpending.gov  -->  Download fiscal-year ZIP  -->  Extract CSV files",
        "   --> Parse & transform to TSV  -->  LOAD DATA INFILE into temp table",
        "   --> ON DUPLICATE KEY UPDATE into usaspending_award",
        "",
        "## Key Details",
        "",
        "Source:           USASpending.gov bulk download (fiscal-year ZIPs)",
        "Format:           CSV files inside ZIP archives",
        "Loading method:   LOAD DATA INFILE (fastest MySQL bulk load, bypasses Python INSERT)",
        "Target table:     usaspending_award (PK: generated_unique_award_id)",
        "Batch size:       50,000 records per LOAD DATA INFILE batch",
        "CLI command:      python main.py load usaspending",
        "",
        "## Column Mapping",
        "",
        "27 CSV columns mapped to DB columns (contract_award_unique_key -> generated_unique_award_id, etc.)",
        "Date columns parsed: start_date, end_date, last_modified_date",
        "Money columns parsed: total_obligation, base_and_all_options_value",
        "",
        ">> Uses TSV intermediate files written to temp dir for LOAD DATA INFILE.",
        ">> Enriches existing award data with FPDS details (fpds_enriched_at timestamp).",
    ])

    # ====================================================================
    # Slide 5: Attachment Intelligence Pipeline
    # ====================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "Attachment Intelligence Pipeline (7-Stage)")
    add_body_text(slide, [
        "## Data Flow (7 Stages)",
        "",
        "1. LOAD:       Load opportunities with resource_links from SAM.gov API",
        "2. ENRICH:     HEAD requests to resolve URLs, get content-type and file size",
        "3. DOWNLOAD:   GET requests to download files to E:\\fedprospector\\attachments\\",
        "4. EXTRACT:    Text extraction (PyMuPDF, python-docx, openpyxl, striprtf)",
        "5. KEYWORD:    Regex/keyword intel scan of extracted Markdown text",
        "6. AI:         AI-powered analysis (Phase 110C, planned)",
        "7. CLEANUP:    Remove downloaded files after all stages complete",
        "",
        "## Text Extraction (Stage 4)",
        "",
        "PDF: PyMuPDF with structure-aware Markdown (headings, bold, tables)",
        "Word: python-docx (.docx), LibreOffice fallback (.doc)",
        "Excel: openpyxl (.xlsx), xlrd (.xls)  |  Also: PPTX, RTF, ODT, CSV, HTML",
        "Scanned PDF detection: < 50 chars/page flagged as image-only",
        "",
        "## Keyword Intel (Stage 5) -- 4 Categories",
        "",
        "Security clearance:  TS/SCI, Top Secret, Secret, Public Trust, SF-86, SCIF",
        "Evaluation method:   LPTA, Best Value, FAR 15.101-1/2, trade-off",
        "Contract vehicle:    OASIS, GSA MAS, BPA, IDIQ, GWAC, SEWP, CIO-SP3/4",
        "Recompete signals:   incumbent, follow-on, re-compete, bridge, transition",
        "",
        ">> SSRF protection on downloads. Redirect validation (*.amazonaws.com only).",
        ">> Parallel downloads (5 workers). Structure-aware confidence boosting from headings.",
    ], font_size=12, line_spacing=1.3)

    # ====================================================================
    # Slide 6: On-Demand Loading
    # ====================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_title_bar(slide, "On-Demand Loading (UI-Triggered)")
    add_body_text(slide, [
        "## Data Flow",
        "",
        "React UI button click  -->  C# API POST /data-load-requests",
        "   --> INSERT into data_load_request (status = PENDING)",
        "   --> Python DemandLoader polls for PENDING rows (up to 10 per cycle)",
        "   --> Route to appropriate loader based on request_type:",
        "       USASPENDING_AWARD  -->  USASpendingClient + USASpendingLoader",
        "       FPDS_AWARD         -->  SAMAwardsClient (key 2) + AwardsLoader",
        "   --> UPDATE data_load_request (status = COMPLETED / FAILED)",
        "   --> UI notification on completion",
        "",
        "## Key Details",
        "",
        "Request table:    data_load_request (PK: request_id)",
        "Status values:    PENDING -> PROCESSING -> COMPLETED | FAILED",
        "Polling:          Python CLI process, periodic poll cycle",
        "API key:          Uses SAM.gov API key 2 (1000/day limit) for FPDS lookups",
        "CLI command:      python main.py demand process",
        "",
        "## Architecture",
        "",
        "Decoupled design: C# API writes request, Python processor reads and executes",
        "Reuses existing ETL loaders -- no duplicate data transformation logic",
        "Error handling: per-request try/catch, failed requests logged with error message",
        "",
        ">> Enables users to trigger targeted data refreshes without running full ETL pipelines.",
    ])

    return prs


if __name__ == "__main__":
    prs = build_presentation()
    output_path = r"c:\git\fedProspect\docs\etl\FedProspect-Data-Flows.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
