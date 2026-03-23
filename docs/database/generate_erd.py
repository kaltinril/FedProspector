"""Generate FedProspect ERD PowerPoint with 7 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette ──────────────────────────────────────────────────────────
BLUE       = RGBColor(0x1E, 0x88, 0xE5)  # Opportunity
BLUE_LIGHT = RGBColor(0xBB, 0xDE, 0xFB)
GREEN      = RGBColor(0x43, 0xA0, 0x47)  # Award
GREEN_LIGHT= RGBColor(0xC8, 0xE6, 0xC9)
ORANGE     = RGBColor(0xFB, 0x8C, 0x00)  # Entity
ORANGE_LIGHT = RGBColor(0xFF, 0xE0, 0xB2)
PURPLE     = RGBColor(0x7B, 0x1F, 0xA2)  # Application
PURPLE_LIGHT = RGBColor(0xE1, 0xBE, 0xE7)
GRAY       = RGBColor(0x61, 0x61, 0x61)  # ETL / Reference
GRAY_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_BG    = RGBColor(0x26, 0x32, 0x38)
TEAL       = RGBColor(0x00, 0x96, 0x88)
TEAL_LIGHT = RGBColor(0xB2, 0xDF, 0xDB)
RED        = RGBColor(0xE5, 0x39, 0x35)
RED_LIGHT  = RGBColor(0xFF, 0xCD, 0xD2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_bar(slide, title_text, subtitle_text=""):
    """Add a dark title bar at the top of the slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BG
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.1)
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)
        p2.alignment = PP_ALIGN.LEFT


def add_table_box(slide, left, top, width, height, table_name, columns,
                  header_color, body_color, pk_col=None, font_size=8):
    """Draw a table entity box with header and column list."""
    # Header
    hdr_h = Inches(0.35)
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_color
    hdr.line.color.rgb = header_color
    hdr.line.width = Pt(1)
    # Round only top corners - adjust
    tf = hdr.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = table_name
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Body
    body_h = height - hdr_h
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + hdr_h, width, body_h)
    body.fill.solid()
    body.fill.fore_color.rgb = body_color
    body.line.color.rgb = header_color
    body.line.width = Pt(1)
    tf2 = body.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.08)
    tf2.margin_right = Inches(0.08)
    tf2.margin_top = Inches(0.04)
    tf2.margin_bottom = Inches(0.04)
    for i, col in enumerate(columns):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        is_pk = (pk_col and col.startswith(pk_col))
        prefix = "PK " if is_pk else "   "
        p.text = prefix + col
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        if is_pk:
            p.font.bold = True
        p.space_after = Pt(1)
        p.space_before = Pt(0)

    return (left, top, width, height)


def add_arrow(slide, x1, y1, x2, y2, label="", color=BLACK):
    """Draw a connector line with optional label."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    if label:
        mid_x = (x1 + x2) // 2
        mid_y = (y1 + y2) // 2
        lbl = slide.shapes.add_textbox(mid_x - Inches(0.5), mid_y - Inches(0.15), Inches(1), Inches(0.3))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(7)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER


def add_cardinality_label(slide, x, y, text, color=BLACK):
    """Add a small cardinality label at a specific position."""
    lbl = slide.shapes.add_textbox(x, y, Inches(0.6), Inches(0.2))
    tf = lbl.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER


def add_domain_box(slide, left, top, width, height, label, fill_color, border_color, tables_text=""):
    """Add a colored domain box for the overview slide."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = border_color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = border_color
    p.alignment = PP_ALIGN.CENTER
    if tables_text:
        p2 = tf.add_paragraph()
        p2.text = tables_text
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0x42, 0x42, 0x42)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
    return box


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1: High-Level Domain Overview
# ══════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_title_bar(slide1, "FedProspect — Entity Relationship Diagram",
              "High-Level Domain Overview  |  fed_contracts database")

# Domain boxes - arranged in a clean layout
# Row 1: Public data domains
opp_box = add_domain_box(slide1, Inches(0.8), Inches(1.3), Inches(3.5), Inches(2.5),
    "Opportunities", BLUE_LIGHT, BLUE,
    "opportunity\nopportunity_history\nopportunity_relationship\nopportunity_attachment\n"
    "opportunity_attachment_intel\nopportunity_intel_source\nopportunity_poc\ncontracting_officer")

award_box = add_domain_box(slide1, Inches(4.9), Inches(1.3), Inches(3.5), Inches(2.5),
    "Awards & Contracts", GREEN_LIGHT, GREEN,
    "usaspending_award\nusaspending_transaction\nfpds_contract\ngsa_labor_rate\n"
    "sam_subaward\nusaspending_load_checkpoint")

entity_box = add_domain_box(slide1, Inches(9.0), Inches(1.3), Inches(3.5), Inches(2.5),
    "Entities (SAM.gov)", ORANGE_LIGHT, ORANGE,
    "entity\nentity_address\nentity_naics\nentity_psc\n"
    "entity_business_type\nentity_sba_certification\nentity_poc\n"
    "entity_disaster_response\nentity_history\nsam_exclusion")

# Row 2: Application + operational
app_box = add_domain_box(slide1, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.8),
    "Application (Multi-Tenant)", PURPLE_LIGHT, PURPLE,
    "organization  |  app_user  |  app_session  |  organization_invite\n"
    "organization_naics  |  organization_certification  |  organization_past_performance\n"
    "organization_entity  |  prospect  |  prospect_note  |  prospect_team_member\n"
    "saved_search  |  notification  |  activity_log\n"
    "proposal  |  proposal_document  |  proposal_milestone")

ref_box = add_domain_box(slide1, Inches(6.8), Inches(4.3), Inches(2.8), Inches(2.8),
    "Reference Data", GRAY_LIGHT, GRAY,
    "ref_naics_code\nref_sba_size_standard\nref_naics_footnote\nref_psc_code\n"
    "ref_set_aside_type\nref_sba_type\nref_business_type\nref_entity_structure\n"
    "ref_country_code  |  ref_state_code\nref_fips_county")

etl_box = add_domain_box(slide1, Inches(10.1), Inches(4.3), Inches(2.4), Inches(2.8),
    "ETL Operational", GRAY_LIGHT, GRAY,
    "etl_load_log\netl_load_error\netl_data_quality_rule\netl_rate_limit\n"
    "etl_health_snapshot\ndata_load_request\nfederal_organization\nstg_entity_raw")

# Arrows between domains
# Opportunities <-> Awards (solicitation_identifier, notice_id linkage)
add_arrow(slide1, Inches(4.3), Inches(2.5), Inches(4.9), Inches(2.5), "1:N", BLUE)
# Awards <-> Entities (recipient_uei, vendor_uei)
add_arrow(slide1, Inches(8.4), Inches(2.5), Inches(9.0), Inches(2.5), "N:1", GREEN)
# Opportunities -> Application (prospect.notice_id)
add_arrow(slide1, Inches(2.5), Inches(3.8), Inches(2.5), Inches(4.3), "1:N", PURPLE)
# Entities -> Application (organization_entity.uei_sam)
add_arrow(slide1, Inches(10.7), Inches(3.8), Inches(5.5), Inches(4.3), "N:M", PURPLE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2: Opportunity Domain
# ══════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide2, "Opportunity Domain", "SAM.gov contract opportunities and attachment intelligence")

# opportunity (center-left)
add_table_box(slide2, Inches(0.3), Inches(1.2), Inches(2.8), Inches(4.5),
    "opportunity", [
        "notice_id VARCHAR(100)",
        "title VARCHAR(500)",
        "solicitation_number VARCHAR(100)",
        "department_name VARCHAR(200)",
        "sub_tier / office",
        "posted_date DATE",
        "response_deadline DATETIME",
        "type / base_type VARCHAR(50)",
        "set_aside_code VARCHAR(20)",
        "naics_code VARCHAR(6)",
        "classification_code VARCHAR(10)",
        "pop_state / pop_zip / pop_country",
        "award_number / award_date",
        "award_amount DECIMAL(15,2)",
        "awardee_uei / awardee_name",
        "description_url / description_text",
        "active CHAR(1)",
        "record_hash CHAR(64)",
        "first_loaded_at / last_loaded_at",
    ], BLUE, BLUE_LIGHT, pk_col="notice_id")

# opportunity_history (top-right of opportunity)
add_table_box(slide2, Inches(3.5), Inches(1.2), Inches(2.3), Inches(1.8),
    "opportunity_history", [
        "id BIGINT AUTO_INCREMENT",
        "notice_id VARCHAR(100) FK",
        "field_name VARCHAR(100)",
        "old_value TEXT",
        "new_value TEXT",
        "changed_at DATETIME",
        "load_id INT",
    ], BLUE, BLUE_LIGHT, pk_col="id")
add_arrow(slide2, Inches(3.1), Inches(2.0), Inches(3.5), Inches(2.0), "1:N", BLUE)

# opportunity_relationship
add_table_box(slide2, Inches(3.5), Inches(3.3), Inches(2.3), Inches(1.6),
    "opportunity_relationship", [
        "id INT AUTO_INCREMENT",
        "parent_notice_id FK",
        "child_notice_id FK",
        "relationship_type VARCHAR(30)",
        "created_by INT",
        "notes TEXT",
    ], BLUE, BLUE_LIGHT, pk_col="id")
add_arrow(slide2, Inches(3.1), Inches(3.5), Inches(3.5), Inches(3.8), "1:N", BLUE)

# opportunity_attachment (center)
add_table_box(slide2, Inches(6.2), Inches(1.2), Inches(2.6), Inches(2.8),
    "opportunity_attachment", [
        "attachment_id INT AUTO_INCREMENT",
        "notice_id VARCHAR(100) FK",
        "url VARCHAR(500)",
        "filename VARCHAR(500)",
        "content_type VARCHAR(100)",
        "file_size_bytes BIGINT",
        "file_path VARCHAR(500)",
        "extracted_text LONGTEXT",
        "download_status ENUM",
        "extraction_status ENUM",
        "content_hash / text_hash",
    ], BLUE, BLUE_LIGHT, pk_col="attachment_id")
add_arrow(slide2, Inches(3.1), Inches(2.8), Inches(6.2), Inches(2.5), "1:N", BLUE)
add_cardinality_label(slide2, Inches(3.1), Inches(2.6), "1", BLUE)
add_cardinality_label(slide2, Inches(5.8), Inches(2.3), "N", BLUE)

# opportunity_attachment_intel
add_table_box(slide2, Inches(9.2), Inches(1.2), Inches(3.0), Inches(3.2),
    "opportunity_attachment_intel", [
        "intel_id INT AUTO_INCREMENT",
        "notice_id VARCHAR(100) FK",
        "attachment_id INT FK",
        "extraction_method ENUM",
        "clearance_required CHAR(1)",
        "clearance_level / clearance_scope",
        "eval_method / eval_details",
        "vehicle_type / vehicle_details",
        "is_recompete CHAR(1)",
        "incumbent_name VARCHAR(200)",
        "scope_summary TEXT",
        "labor_categories JSON",
        "key_requirements JSON",
        "overall_confidence ENUM",
    ], BLUE, BLUE_LIGHT, pk_col="intel_id")
add_arrow(slide2, Inches(8.8), Inches(2.5), Inches(9.2), Inches(2.5), "1:N", BLUE)

# opportunity_intel_source
add_table_box(slide2, Inches(9.2), Inches(4.7), Inches(3.0), Inches(2.3),
    "opportunity_intel_source", [
        "source_id INT AUTO_INCREMENT",
        "intel_id INT FK",
        "field_name VARCHAR(50)",
        "attachment_id INT FK",
        "source_filename VARCHAR(500)",
        "page_number INT",
        "matched_text VARCHAR(500)",
        "pattern_name VARCHAR(100)",
        "extraction_method ENUM",
        "confidence ENUM",
    ], BLUE, BLUE_LIGHT, pk_col="source_id")
add_arrow(slide2, Inches(10.5), Inches(4.4), Inches(10.5), Inches(4.7), "1:N", BLUE)
add_cardinality_label(slide2, Inches(10.2), Inches(4.2), "1", BLUE)
add_cardinality_label(slide2, Inches(10.2), Inches(4.7), "N", BLUE)

# opportunity_poc & contracting_officer
add_table_box(slide2, Inches(3.5), Inches(5.2), Inches(2.3), Inches(1.2),
    "opportunity_poc", [
        "poc_id INT AUTO_INCREMENT",
        "notice_id VARCHAR(100) FK",
        "officer_id INT FK",
        "poc_type VARCHAR(20)",
    ], BLUE, BLUE_LIGHT, pk_col="poc_id")
add_arrow(slide2, Inches(2.5), Inches(5.7), Inches(3.5), Inches(5.7), "1:N", BLUE)

add_table_box(slide2, Inches(6.2), Inches(5.0), Inches(2.6), Inches(1.8),
    "contracting_officer", [
        "officer_id INT AUTO_INCREMENT",
        "full_name VARCHAR(500)",
        "email / phone / fax",
        "title VARCHAR(200)",
        "department_name VARCHAR(200)",
        "office_name VARCHAR(200)",
    ], BLUE, BLUE_LIGHT, pk_col="officer_id")
add_arrow(slide2, Inches(5.8), Inches(5.7), Inches(6.2), Inches(5.7), "N:1", BLUE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Award Domain
# ══════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide3, "Award & Contract Domain",
              "USASpending.gov awards, FPDS contracts, GSA rates, subawards")

# usaspending_award (left)
add_table_box(slide3, Inches(0.3), Inches(1.2), Inches(3.0), Inches(4.2),
    "usaspending_award", [
        "generated_unique_award_id PK",
        "piid / fain / uri",
        "award_type VARCHAR(50)",
        "award_description TEXT",
        "recipient_name / recipient_uei",
        "recipient_parent_name / parent_uei",
        "total_obligation DECIMAL(15,2)",
        "base_and_all_options_value",
        "start_date / end_date",
        "awarding_agency_name",
        "funding_agency_name",
        "naics_code / psc_code",
        "type_of_set_aside",
        "pop_state / pop_country / pop_zip",
        "solicitation_identifier",
        "fiscal_year SMALLINT",
        "record_hash CHAR(64)",
    ], GREEN, GREEN_LIGHT, pk_col="generated_unique_award_id")

# usaspending_transaction
add_table_box(slide3, Inches(3.7), Inches(1.2), Inches(2.6), Inches(2.0),
    "usaspending_transaction", [
        "id BIGINT AUTO_INCREMENT",
        "award_id VARCHAR(100) FK",
        "action_date DATE",
        "modification_number VARCHAR(25)",
        "action_type VARCHAR(5)",
        "federal_action_obligation DECIMAL",
        "description TEXT",
    ], GREEN, GREEN_LIGHT, pk_col="id")
add_arrow(slide3, Inches(3.3), Inches(2.2), Inches(3.7), Inches(2.2), "1:N", GREEN)
add_cardinality_label(slide3, Inches(3.0), Inches(2.0), "1", GREEN)
add_cardinality_label(slide3, Inches(3.7), Inches(2.0), "N", GREEN)

# usaspending_load_checkpoint
add_table_box(slide3, Inches(3.7), Inches(3.5), Inches(2.6), Inches(1.8),
    "usaspending_load_checkpoint", [
        "checkpoint_id INT AUTO_INCREMENT",
        "load_id INT FK (etl_load_log)",
        "fiscal_year INT",
        "csv_file_name VARCHAR(255)",
        "status ENUM",
        "completed_batches INT",
        "total_rows_loaded INT",
        "archive_hash VARCHAR(130)",
    ], GREEN, GREEN_LIGHT, pk_col="checkpoint_id")

# fpds_contract (center-right)
add_table_box(slide3, Inches(6.8), Inches(1.2), Inches(2.8), Inches(4.5),
    "fpds_contract", [
        "contract_id + modification_number PK",
        "idv_piid VARCHAR(50)",
        "agency_id / agency_name",
        "contracting_office_id / name",
        "funding_agency_id / name",
        "vendor_uei / vendor_name",
        "date_signed / effective_date",
        "completion_date / ultimate_compl.",
        "dollars_obligated DECIMAL(15,2)",
        "base_and_all_options DECIMAL",
        "naics_code / psc_code",
        "set_aside_type VARCHAR(20)",
        "type_of_contract VARCHAR(10)",
        "description TEXT",
        "solicitation_number / date",
        "number_of_offers INT",
        "record_hash CHAR(64)",
    ], GREEN, GREEN_LIGHT, pk_col="contract_id")

# sam_subaward (right)
add_table_box(slide3, Inches(10.0), Inches(1.2), Inches(2.8), Inches(3.2),
    "sam_subaward", [
        "id INT AUTO_INCREMENT",
        "prime_piid VARCHAR(50)",
        "prime_agency_id / prime_agency_name",
        "prime_uei / prime_name",
        "sub_uei / sub_name",
        "sub_amount DECIMAL(15,2)",
        "sub_date DATE",
        "sub_description TEXT",
        "naics_code / psc_code",
        "sub_business_type VARCHAR(50)",
        "pop_state / pop_country / pop_zip",
        "record_hash CHAR(64)",
    ], GREEN, GREEN_LIGHT, pk_col="id")

# gsa_labor_rate
add_table_box(slide3, Inches(10.0), Inches(4.8), Inches(2.8), Inches(2.2),
    "gsa_labor_rate", [
        "id INT AUTO_INCREMENT",
        "labor_category VARCHAR(200)",
        "education_level / min_years_exp",
        "current_price / next_year_price",
        "schedule VARCHAR(200)",
        "contractor_name VARCHAR(200)",
        "sin VARCHAR(500)",
        "business_size VARCHAR(10)",
        "security_clearance VARCHAR(50)",
        "idv_piid VARCHAR(50)",
    ], GREEN, GREEN_LIGHT, pk_col="id")

# Relationship labels
lbl = slide3.shapes.add_textbox(Inches(4.0), Inches(5.6), Inches(5.0), Inches(0.5))
tf = lbl.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Linked by: solicitation_identifier <-> solicitation_number, piid, vendor_uei/recipient_uei"
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = GREEN


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4: Entity Domain
# ══════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide4, "Entity Domain (SAM.gov)",
              "Contractor registrations, certifications, exclusions, and points of contact")

# entity (center)
add_table_box(slide4, Inches(4.5), Inches(1.2), Inches(3.2), Inches(4.0),
    "entity", [
        "uei_sam VARCHAR(12)",
        "cage_code VARCHAR(5)",
        "registration_status CHAR(1)",
        "legal_business_name VARCHAR(120)",
        "dba_name VARCHAR(120)",
        "entity_url VARCHAR(200)",
        "primary_naics VARCHAR(6)",
        "entity_structure_code VARCHAR(2)",
        "state_of_incorporation VARCHAR(2)",
        "country_of_incorporation VARCHAR(3)",
        "exclusion_status_flag CHAR(1)",
        "initial_registration_date DATE",
        "registration_expiration_date DATE",
        "record_hash CHAR(64)",
        "first_loaded_at / last_loaded_at",
    ], ORANGE, ORANGE_LIGHT, pk_col="uei_sam")

# entity_address (top-left)
add_table_box(slide4, Inches(0.3), Inches(1.2), Inches(2.5), Inches(1.6),
    "entity_address", [
        "id INT AUTO_INCREMENT",
        "uei_sam VARCHAR(12) FK",
        "address_type VARCHAR(10)",
        "address_line_1/2",
        "city / state_or_province",
        "zip_code / country_code",
    ], ORANGE, ORANGE_LIGHT, pk_col="id")
add_arrow(slide4, Inches(2.8), Inches(2.0), Inches(4.5), Inches(2.0), "1:N", ORANGE)
add_cardinality_label(slide4, Inches(4.2), Inches(1.8), "1", ORANGE)
add_cardinality_label(slide4, Inches(2.5), Inches(1.8), "N", ORANGE)

# entity_naics (left-mid)
add_table_box(slide4, Inches(0.3), Inches(3.1), Inches(2.5), Inches(1.3),
    "entity_naics", [
        "id INT AUTO_INCREMENT",
        "uei_sam VARCHAR(12) FK",
        "naics_code VARCHAR(11)",
        "is_primary CHAR(1)",
        "sba_small_business CHAR(1)",
    ], ORANGE, ORANGE_LIGHT, pk_col="id")
add_arrow(slide4, Inches(2.8), Inches(3.7), Inches(4.5), Inches(3.2), "1:N", ORANGE)

# entity_psc (left-bottom)
add_table_box(slide4, Inches(0.3), Inches(4.7), Inches(2.5), Inches(1.0),
    "entity_psc", [
        "id INT AUTO_INCREMENT",
        "uei_sam VARCHAR(12) FK",
        "psc_code VARCHAR(10)",
    ], ORANGE, ORANGE_LIGHT, pk_col="id")
add_arrow(slide4, Inches(2.8), Inches(5.1), Inches(4.5), Inches(4.5), "1:N", ORANGE)

# entity_business_type (bottom-left)
add_table_box(slide4, Inches(0.3), Inches(6.0), Inches(2.5), Inches(1.0),
    "entity_business_type", [
        "id INT AUTO_INCREMENT",
        "uei_sam VARCHAR(12) FK",
        "business_type_code VARCHAR(4)",
    ], ORANGE, ORANGE_LIGHT, pk_col="id")
add_arrow(slide4, Inches(2.8), Inches(6.3), Inches(4.5), Inches(5.0), "1:N", ORANGE)

# entity_poc (top-right)
add_table_box(slide4, Inches(8.2), Inches(1.2), Inches(2.8), Inches(2.0),
    "entity_poc", [
        "id INT AUTO_INCREMENT",
        "uei_sam VARCHAR(12) FK",
        "poc_type VARCHAR(40)",
        "first_name / last_name",
        "title VARCHAR(50)",
        "address / city / state / zip",
    ], ORANGE, ORANGE_LIGHT, pk_col="id")
add_arrow(slide4, Inches(7.7), Inches(2.2), Inches(8.2), Inches(2.2), "1:N", ORANGE)

# entity_sba_certification (right-mid)
add_table_box(slide4, Inches(8.2), Inches(3.5), Inches(2.8), Inches(1.5),
    "entity_sba_certification", [
        "id INT AUTO_INCREMENT",
        "uei_sam VARCHAR(12) FK",
        "sba_type_code / sba_type_desc",
        "certification_entry_date",
        "certification_exit_date",
    ], ORANGE, ORANGE_LIGHT, pk_col="id")
add_arrow(slide4, Inches(7.7), Inches(3.8), Inches(8.2), Inches(4.0), "1:N", ORANGE)

# entity_disaster_response (right-bottom)
add_table_box(slide4, Inches(8.2), Inches(5.3), Inches(2.8), Inches(1.3),
    "entity_disaster_response", [
        "id INT AUTO_INCREMENT",
        "uei_sam VARCHAR(12) FK",
        "state_code / state_name",
        "county_code / county_name",
        "msa_code / msa_name",
    ], ORANGE, ORANGE_LIGHT, pk_col="id")
add_arrow(slide4, Inches(7.7), Inches(5.0), Inches(8.2), Inches(5.8), "1:N", ORANGE)

# entity_history (bottom-right)
add_table_box(slide4, Inches(11.3), Inches(1.2), Inches(1.7), Inches(1.8),
    "entity_history", [
        "id BIGINT AUTO_INC",
        "uei_sam FK",
        "field_name",
        "old_value TEXT",
        "new_value TEXT",
        "changed_at",
        "load_id INT",
    ], ORANGE, ORANGE_LIGHT, pk_col="id", font_size=7)

# sam_exclusion (bottom-center)
add_table_box(slide4, Inches(4.5), Inches(5.5), Inches(3.2), Inches(1.8),
    "sam_exclusion", [
        "id INT AUTO_INCREMENT",
        "uei VARCHAR(12)",
        "cage_code VARCHAR(10)",
        "entity_name / first_name / last_name",
        "exclusion_type / exclusion_program",
        "excluding_agency_code / name",
        "activation_date / termination_date",
        "record_hash CHAR(64)",
    ], RED, RED_LIGHT, pk_col="id")
# Note: sam_exclusion links by uei but no FK constraint
lbl = slide4.shapes.add_textbox(Inches(4.5), Inches(7.3), Inches(4.0), Inches(0.3))
tf = lbl.text_frame
p = tf.paragraphs[0]
p.text = "sam_exclusion links to entity by uei (logical, no FK constraint)"
p.font.size = Pt(8)
p.font.italic = True
p.font.color.rgb = RED


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5: Application Domain
# ══════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide5, "Application Domain (Multi-Tenant)",
              "Organizations, users, prospects, proposals, and saved searches")

# organization (top-left)
add_table_box(slide5, Inches(0.3), Inches(1.2), Inches(2.8), Inches(3.0),
    "organization", [
        "organization_id INT AUTO_INCREMENT",
        "name / slug VARCHAR",
        "is_active CHAR(1)",
        "max_users INT",
        "subscription_tier VARCHAR(50)",
        "legal_name / dba_name",
        "uei_sam / cage_code / ein",
        "phone / website",
        "address / city / state / zip",
        "annual_revenue DECIMAL(18,2)",
        "employee_count INT",
        "profile_completed CHAR(1)",
    ], PURPLE, PURPLE_LIGHT, pk_col="organization_id")

# app_user (top-center)
add_table_box(slide5, Inches(3.5), Inches(1.2), Inches(2.5), Inches(2.5),
    "app_user", [
        "user_id INT AUTO_INCREMENT",
        "organization_id INT FK",
        "username / display_name",
        "email / password_hash",
        "role VARCHAR(20)",
        "is_org_admin / org_role",
        "is_system_admin TINYINT",
        "mfa_enabled CHAR(1)",
        "last_login_at DATETIME",
        "is_active CHAR(1)",
    ], PURPLE, PURPLE_LIGHT, pk_col="user_id")
add_arrow(slide5, Inches(3.1), Inches(2.5), Inches(3.5), Inches(2.5), "1:N", PURPLE)
add_cardinality_label(slide5, Inches(2.8), Inches(2.3), "1", PURPLE)
add_cardinality_label(slide5, Inches(3.5), Inches(2.3), "N", PURPLE)

# prospect (center)
add_table_box(slide5, Inches(6.4), Inches(1.2), Inches(2.8), Inches(3.2),
    "prospect", [
        "prospect_id INT AUTO_INCREMENT",
        "organization_id INT FK",
        "notice_id VARCHAR(100) FK",
        "assigned_to INT FK (app_user)",
        "capture_manager_id INT FK",
        "status VARCHAR(30)",
        "priority / decision_date",
        "estimated_value DECIMAL(15,2)",
        "win_probability DECIMAL(5,2)",
        "go_no_go_score DECIMAL(5,2)",
        "teaming_required CHAR(1)",
        "outcome / outcome_date",
    ], PURPLE, PURPLE_LIGHT, pk_col="prospect_id")
add_arrow(slide5, Inches(6.0), Inches(2.5), Inches(6.4), Inches(2.5), "1:N", PURPLE)
add_cardinality_label(slide5, Inches(5.7), Inches(2.3), "1", PURPLE)
add_cardinality_label(slide5, Inches(6.4), Inches(2.3), "N", PURPLE)

# prospect_note (right of prospect)
add_table_box(slide5, Inches(9.6), Inches(1.2), Inches(2.2), Inches(1.3),
    "prospect_note", [
        "note_id INT AUTO_INCREMENT",
        "prospect_id INT FK",
        "user_id INT FK",
        "note_type VARCHAR(30)",
        "note_text TEXT",
    ], PURPLE, PURPLE_LIGHT, pk_col="note_id")
add_arrow(slide5, Inches(9.2), Inches(1.8), Inches(9.6), Inches(1.8), "1:N", PURPLE)

# prospect_team_member
add_table_box(slide5, Inches(9.6), Inches(2.8), Inches(2.2), Inches(1.4),
    "prospect_team_member", [
        "id INT AUTO_INCREMENT",
        "prospect_id INT FK",
        "uei_sam VARCHAR(12)",
        "app_user_id INT FK",
        "role / notes",
        "proposed_hourly_rate",
    ], PURPLE, PURPLE_LIGHT, pk_col="id")
add_arrow(slide5, Inches(9.2), Inches(3.3), Inches(9.6), Inches(3.3), "1:N", PURPLE)

# proposal (below prospect)
add_table_box(slide5, Inches(6.4), Inches(4.8), Inches(2.8), Inches(1.6),
    "proposal", [
        "proposal_id INT AUTO_INCREMENT",
        "prospect_id INT FK (1:1)",
        "proposal_number / status",
        "submission_deadline DATETIME",
        "estimated_value / win_prob",
        "lessons_learned TEXT",
    ], PURPLE, PURPLE_LIGHT, pk_col="proposal_id")
add_arrow(slide5, Inches(7.8), Inches(4.4), Inches(7.8), Inches(4.8), "1:1", PURPLE)

# proposal_document + proposal_milestone (bottom-right)
add_table_box(slide5, Inches(9.6), Inches(4.5), Inches(2.2), Inches(1.2),
    "proposal_document", [
        "document_id INT",
        "proposal_id INT FK",
        "document_type / file_name",
        "uploaded_by INT FK",
    ], PURPLE, PURPLE_LIGHT, pk_col="document_id")
add_arrow(slide5, Inches(9.2), Inches(5.1), Inches(9.6), Inches(5.1), "1:N", PURPLE)

add_table_box(slide5, Inches(9.6), Inches(5.9), Inches(2.2), Inches(1.2),
    "proposal_milestone", [
        "milestone_id INT",
        "proposal_id INT FK",
        "milestone_name / due_date",
        "assigned_to INT FK",
        "status VARCHAR(20)",
    ], PURPLE, PURPLE_LIGHT, pk_col="milestone_id")
add_arrow(slide5, Inches(9.2), Inches(6.3), Inches(9.6), Inches(6.3), "1:N", PURPLE)

# Org child tables (bottom-left)
add_table_box(slide5, Inches(0.3), Inches(4.5), Inches(2.0), Inches(0.9),
    "organization_naics", [
        "id | organization_id FK",
        "naics_code | is_primary",
        "size_standard_met",
    ], PURPLE, PURPLE_LIGHT, pk_col="id", font_size=7)
add_arrow(slide5, Inches(1.3), Inches(4.2), Inches(1.3), Inches(4.5), "1:N", PURPLE)

add_table_box(slide5, Inches(0.3), Inches(5.6), Inches(2.0), Inches(0.9),
    "organization_certification", [
        "id | organization_id FK",
        "certification_type",
        "expiration_date",
    ], PURPLE, PURPLE_LIGHT, pk_col="id", font_size=7)

add_table_box(slide5, Inches(0.3), Inches(6.7), Inches(2.0), Inches(0.6),
    "organization_past_perf", [
        "id | org_id FK | contract_no",
    ], PURPLE, PURPLE_LIGHT, pk_col="id", font_size=7)

add_table_box(slide5, Inches(2.5), Inches(4.5), Inches(2.0), Inches(0.9),
    "organization_entity", [
        "id | organization_id FK",
        "uei_sam FK | relationship",
        "partner_uei | is_active",
    ], PURPLE, PURPLE_LIGHT, pk_col="id", font_size=7)

# saved_search, notification, activity_log, app_session, org_invite
add_table_box(slide5, Inches(3.5), Inches(4.0), Inches(2.5), Inches(1.2),
    "saved_search", [
        "search_id INT AUTO_INCREMENT",
        "organization_id / user_id FK",
        "search_name / filter_criteria JSON",
        "notification_enabled CHAR(1)",
        "auto_prospect_enabled",
    ], PURPLE, PURPLE_LIGHT, pk_col="search_id")

add_table_box(slide5, Inches(2.5), Inches(5.6), Inches(2.0), Inches(0.8),
    "notification", [
        "notification_id | user_id FK",
        "type / title / is_read",
    ], PURPLE, PURPLE_LIGHT, pk_col="notification_id", font_size=7)

add_table_box(slide5, Inches(2.5), Inches(6.6), Inches(2.0), Inches(0.8),
    "activity_log", [
        "activity_id | org_id / user_id",
        "action / entity_type / details",
    ], PURPLE, PURPLE_LIGHT, pk_col="activity_id", font_size=7)

add_table_box(slide5, Inches(4.8), Inches(5.5), Inches(2.0), Inches(0.8),
    "app_session", [
        "session_id | user_id FK",
        "token_hash | expires_at",
    ], PURPLE, PURPLE_LIGHT, pk_col="session_id", font_size=7)

add_table_box(slide5, Inches(4.8), Inches(6.5), Inches(2.0), Inches(0.8),
    "organization_invite", [
        "invite_id | org_id FK",
        "email | invite_code | org_role",
    ], PURPLE, PURPLE_LIGHT, pk_col="invite_id", font_size=7)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6: ETL Tables
# ══════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide6, "ETL & Operational Tables", "Load tracking, error logging, data quality, rate limiting")

add_table_box(slide6, Inches(0.5), Inches(1.3), Inches(3.0), Inches(3.0),
    "etl_load_log", [
        "load_id INT AUTO_INCREMENT",
        "source_system VARCHAR(50)",
        "load_type VARCHAR(20)",
        "status VARCHAR(20)",
        "started_at / completed_at",
        "records_read INT",
        "records_inserted INT",
        "records_updated INT",
        "records_unchanged INT",
        "records_errored INT",
        "error_message TEXT",
        "parameters JSON",
        "source_file VARCHAR(500)",
    ], GRAY, GRAY_LIGHT, pk_col="load_id")

add_table_box(slide6, Inches(3.9), Inches(1.3), Inches(2.8), Inches(1.8),
    "etl_load_error", [
        "id BIGINT AUTO_INCREMENT",
        "load_id INT FK (etl_load_log)",
        "record_identifier VARCHAR(100)",
        "error_type VARCHAR(50)",
        "error_message TEXT",
        "raw_data TEXT",
        "created_at DATETIME",
    ], GRAY, GRAY_LIGHT, pk_col="id")
add_arrow(slide6, Inches(3.5), Inches(2.2), Inches(3.9), Inches(2.2), "1:N", GRAY)
add_cardinality_label(slide6, Inches(3.2), Inches(2.0), "1", GRAY)
add_cardinality_label(slide6, Inches(3.9), Inches(2.0), "N", GRAY)

add_table_box(slide6, Inches(3.9), Inches(3.4), Inches(2.8), Inches(1.6),
    "etl_data_quality_rule", [
        "rule_id INT AUTO_INCREMENT",
        "rule_name VARCHAR(100)",
        "description TEXT",
        "target_table / target_column",
        "rule_type VARCHAR(20)",
        "rule_definition JSON",
        "is_active / priority",
    ], GRAY, GRAY_LIGHT, pk_col="rule_id")

add_table_box(slide6, Inches(7.2), Inches(1.3), Inches(2.8), Inches(1.5),
    "etl_rate_limit", [
        "id INT AUTO_INCREMENT",
        "source_system VARCHAR(50)",
        "request_date DATE",
        "requests_made / max_requests",
        "last_request_at DATETIME",
    ], GRAY, GRAY_LIGHT, pk_col="id")

add_table_box(slide6, Inches(7.2), Inches(3.1), Inches(2.8), Inches(1.5),
    "etl_health_snapshot", [
        "snapshot_id INT AUTO_INCREMENT",
        "checked_at DATETIME",
        "overall_status VARCHAR(20)",
        "results_json JSON",
        "alert_count / error_count",
        "stale_source_count INT",
    ], GRAY, GRAY_LIGHT, pk_col="snapshot_id")

add_table_box(slide6, Inches(7.2), Inches(4.9), Inches(2.8), Inches(1.8),
    "data_load_request", [
        "request_id INT AUTO_INCREMENT",
        "request_type VARCHAR(30)",
        "lookup_key / lookup_key_type",
        "status VARCHAR(20)",
        "requested_by INT",
        "started_at / completed_at",
        "load_id INT",
        "error_message / result_summary",
    ], GRAY, GRAY_LIGHT, pk_col="request_id")

add_table_box(slide6, Inches(0.5), Inches(4.6), Inches(3.0), Inches(2.5),
    "federal_organization", [
        "fh_org_id INT",
        "fh_org_name VARCHAR(500)",
        "fh_org_type VARCHAR(50)",
        "description TEXT",
        "status VARCHAR(20)",
        "agency_code / cgac",
        "parent_org_id INT (self-ref)",
        "level INT",
        "created_date / last_modified_date",
        "record_hash CHAR(64)",
    ], TEAL, TEAL_LIGHT, pk_col="fh_org_id")

add_table_box(slide6, Inches(10.3), Inches(1.3), Inches(2.5), Inches(2.0),
    "stg_entity_raw", [
        "load_id INT",
        "uei_sam VARCHAR(12)",
        "raw_json JSON",
        "raw_record_hash CHAR(64)",
        "processed CHAR(1)",
        "processed_at DATETIME",
        "error_message TEXT",
    ], GRAY, GRAY_LIGHT, pk_col="load_id")

# Legend
legend = slide6.shapes.add_textbox(Inches(10.3), Inches(4.0), Inches(2.5), Inches(1.5))
tf = legend.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Referenced by all domain tables:"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = GRAY
p2 = tf.add_paragraph()
p2.text = ("Every data table has last_load_id "
           "linking back to etl_load_log.load_id "
           "for load provenance tracking.")
p2.font.size = Pt(8)
p2.font.color.rgb = GRAY


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7: Reference Tables
# ══════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide7, "Reference / Lookup Tables",
              "Loaded from CSV files, updated infrequently  |  Prefixed with ref_")

# Row 1
add_table_box(slide7, Inches(0.3), Inches(1.3), Inches(2.8), Inches(2.2),
    "ref_naics_code", [
        "naics_code VARCHAR(11)",
        "description VARCHAR(500)",
        "code_level TINYINT",
        "level_name VARCHAR(30)",
        "parent_code VARCHAR(11)",
        "year_version VARCHAR(4)",
        "is_active CHAR(1)",
        "footnote_id VARCHAR(5)",
    ], GRAY, GRAY_LIGHT, pk_col="naics_code")

add_table_box(slide7, Inches(3.4), Inches(1.3), Inches(2.6), Inches(1.6),
    "ref_sba_size_standard", [
        "id INT AUTO_INCREMENT",
        "naics_code VARCHAR(11) FK",
        "industry_description",
        "size_standard DECIMAL(13,2)",
        "size_type CHAR(1)",
        "effective_date DATE",
    ], GRAY, GRAY_LIGHT, pk_col="id")
add_arrow(slide7, Inches(3.1), Inches(2.1), Inches(3.4), Inches(2.1), "1:N", GRAY)
add_cardinality_label(slide7, Inches(2.8), Inches(1.9), "1", GRAY)
add_cardinality_label(slide7, Inches(3.4), Inches(1.9), "N", GRAY)

add_table_box(slide7, Inches(3.4), Inches(3.2), Inches(2.6), Inches(1.0),
    "ref_naics_footnote", [
        "footnote_id + section PK",
        "description TEXT",
    ], GRAY, GRAY_LIGHT, pk_col="footnote_id")

add_table_box(slide7, Inches(6.3), Inches(1.3), Inches(3.0), Inches(2.5),
    "ref_psc_code", [
        "psc_code + start_date PK",
        "psc_name VARCHAR(200)",
        "end_date DATE",
        "full_description TEXT",
        "psc_includes / psc_excludes",
        "psc_notes TEXT",
        "parent_psc_code VARCHAR(200)",
        "category_type CHAR(1)",
        "level1_category_code / name",
        "level2_category_code / name",
    ], GRAY, GRAY_LIGHT, pk_col="psc_code")

add_table_box(slide7, Inches(9.7), Inches(1.3), Inches(3.0), Inches(1.5),
    "ref_set_aside_type", [
        "set_aside_code VARCHAR(10)",
        "description VARCHAR(200)",
        "is_small_business CHAR(1)",
        "category VARCHAR(50)",
    ], GRAY, GRAY_LIGHT, pk_col="set_aside_code")

add_table_box(slide7, Inches(9.7), Inches(3.0), Inches(3.0), Inches(1.2),
    "ref_sba_type", [
        "sba_type_code VARCHAR(10)",
        "description VARCHAR(200)",
        "program_name VARCHAR(100)",
    ], GRAY, GRAY_LIGHT, pk_col="sba_type_code")

# Row 2
add_table_box(slide7, Inches(0.3), Inches(4.0), Inches(2.8), Inches(1.5),
    "ref_business_type", [
        "business_type_code VARCHAR(4)",
        "description VARCHAR(200)",
        "classification / category",
        "is_socioeconomic CHAR(1)",
        "is_small_business_related",
    ], GRAY, GRAY_LIGHT, pk_col="business_type_code")

add_table_box(slide7, Inches(0.3), Inches(5.8), Inches(2.8), Inches(1.0),
    "ref_entity_structure", [
        "structure_code VARCHAR(2)",
        "description VARCHAR(200)",
    ], GRAY, GRAY_LIGHT, pk_col="structure_code")

add_table_box(slide7, Inches(3.4), Inches(4.5), Inches(2.6), Inches(1.5),
    "ref_country_code", [
        "three_code VARCHAR(3)",
        "country_name VARCHAR(100)",
        "two_code VARCHAR(2)",
        "numeric_code VARCHAR(4)",
        "is_iso_standard CHAR(1)",
        "sam_gov_recognized CHAR(1)",
    ], GRAY, GRAY_LIGHT, pk_col="three_code")

add_table_box(slide7, Inches(6.3), Inches(4.2), Inches(2.5), Inches(1.0),
    "ref_state_code", [
        "state_code + country_code PK",
        "state_name VARCHAR(60)",
    ], GRAY, GRAY_LIGHT, pk_col="state_code")

add_table_box(slide7, Inches(6.3), Inches(5.5), Inches(2.5), Inches(1.0),
    "ref_fips_county", [
        "fips_code VARCHAR(5)",
        "county_name VARCHAR(100)",
        "state_name VARCHAR(60)",
    ], GRAY, GRAY_LIGHT, pk_col="fips_code")

# Usage note
note = slide7.shapes.add_textbox(Inches(9.7), Inches(4.5), Inches(3.0), Inches(2.5))
tf = note.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Referenced by:"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = GRAY
refs = [
    ("ref_naics_code", "opportunity.naics_code, entity.primary_naics, entity_naics.naics_code, usaspending_award.naics_code, fpds_contract.naics_code"),
    ("ref_psc_code", "entity_psc.psc_code, usaspending_award.psc_code, fpds_contract.psc_code"),
    ("ref_set_aside_type", "opportunity.set_aside_code"),
    ("ref_sba_type", "entity_sba_certification.sba_type_code"),
    ("ref_business_type", "entity_business_type.business_type_code"),
]
for name, usage in refs:
    p2 = tf.add_paragraph()
    p2.text = f"{name}"
    p2.font.size = Pt(8)
    p2.font.bold = True
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(4)
    p3 = tf.add_paragraph()
    p3.text = f"  -> {usage}"
    p3.font.size = Pt(7)
    p3.font.color.rgb = RGBColor(0x75, 0x75, 0x75)


# ── Save ──────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FedProspect-ERD.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")
