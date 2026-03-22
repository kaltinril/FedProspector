"""Manual tests for Phase 110B document type extractors.

Tests: _extract_pptx, _extract_xls, _extract_rtf, _extract_odt, _extract_doc
Each test creates a small file programmatically, runs the extractor, and verifies output.
"""

import os
import sys
import tempfile
import traceback

# Ensure fed_prospector is on the path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from etl.attachment_text_extractor import AttachmentTextExtractor

extractor = AttachmentTextExtractor()
results = []


def report(name, passed, detail=""):
    status = "PASS" if passed else "FAIL"
    results.append((name, passed))
    print(f"\n{'='*60}")
    print(f"[{status}] {name}")
    if detail:
        print(detail)
    print(f"{'='*60}")


# =====================================================================
# Test 1: .pptx
# =====================================================================
def test_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")

        prs = Presentation()
        # Slide 1: title + body + notes
        slide1_layout = prs.slide_layouts[1]  # Title and Content
        slide1 = prs.slides.add_slide(slide1_layout)
        slide1.shapes.title.text = "Test Title"
        slide1.placeholders[1].text = "Body content"
        # Add speaker notes
        notes_slide = slide1.notes_slide
        notes_slide.notes_text_frame.text = "These are notes"

        # Slide 2: table 2x2
        slide2_layout = prs.slide_layouts[5]  # Blank
        slide2 = prs.slides.add_slide(slide2_layout)
        table_shape = slide2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "Col A"
        table.cell(0, 1).text = "Col B"
        table.cell(1, 0).text = "Val 1"
        table.cell(1, 1).text = "Val 2"

        prs.save(path)

        text, page_count, is_scanned = extractor._extract_pptx(path)

        print(f"\n--- Extracted PPTX text ---\n{text}\n--- end ---")

        checks = []
        checks.append(("'## Slide 1' in text", "## Slide 1" in text))
        checks.append(("'## Slide 2' in text", "## Slide 2" in text))
        checks.append(("'Test Title' in text", "Test Title" in text))
        checks.append(("'Body content' in text", "Body content" in text))
        checks.append(("'Notes:' in text", "Notes:" in text))
        checks.append(("'These are notes' in text", "These are notes" in text))
        checks.append(("page_count == 2", page_count == 2))
        checks.append(("is_scanned == False", is_scanned == False))
        # Table content
        checks.append(("'Col A' in text", "Col A" in text))
        checks.append(("'Val 1' in text", "Val 1" in text))

        failed = [c for c in checks if not c[1]]
        detail = "\n".join(f"  {'OK' if c[1] else 'FAIL'}: {c[0]}" for c in checks)
        report("PPTX Extraction", len(failed) == 0, detail)


# =====================================================================
# Test 2: .xls
# =====================================================================
def test_xls():
    import xlwt

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.xls")

        wb = xlwt.Workbook()
        ws = wb.add_sheet("Pricing")
        ws.write(0, 0, "Item")
        ws.write(0, 1, "Price")
        ws.write(1, 0, "Widget")
        ws.write(1, 1, 9.99)
        wb.save(path)

        text, page_count, is_scanned = extractor._extract_xls(path)

        print(f"\n--- Extracted XLS text ---\n{text}\n--- end ---")

        checks = []
        checks.append(("'## Pricing' in text", "## Pricing" in text))
        checks.append(("'Item' in text", "Item" in text))
        checks.append(("'Price' in text", "Price" in text))
        checks.append(("'Widget' in text", "Widget" in text))
        checks.append(("'9.99' in text", "9.99" in text))
        checks.append(("page_count == 0", page_count == 0))
        checks.append(("is_scanned == False", is_scanned == False))

        failed = [c for c in checks if not c[1]]
        detail = "\n".join(f"  {'OK' if c[1] else 'FAIL'}: {c[0]}" for c in checks)
        report("XLS Extraction", len(failed) == 0, detail)


# =====================================================================
# Test 3: .rtf
# =====================================================================
def test_rtf():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.rtf")

        rtf_content = r'{\rtf1\ansi{\fonttbl{\f0 Times New Roman;}}\f0\fs24 This is a test RTF document with some content.}'
        with open(path, "w", encoding="utf-8") as f:
            f.write(rtf_content)

        text, page_count, is_scanned = extractor._extract_rtf(path)

        print(f"\n--- Extracted RTF text ---\n{text}\n--- end ---")

        checks = []
        checks.append(("'This is a test RTF document' in text", "This is a test RTF document" in text))
        checks.append(("page_count == 0", page_count == 0))
        checks.append(("is_scanned == False", is_scanned == False))

        failed = [c for c in checks if not c[1]]
        detail = "\n".join(f"  {'OK' if c[1] else 'FAIL'}: {c[0]}" for c in checks)
        report("RTF Extraction", len(failed) == 0, detail)


# =====================================================================
# Test 4: .odt
# =====================================================================
def test_odt():
    from odf.opendocument import OpenDocumentText
    from odf.text import H, P

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.odt")

        doc = OpenDocumentText()
        heading = H(outlinelevel="1", text="Introduction")
        doc.text.addElement(heading)
        para = P(text="This is the body text.")
        doc.text.addElement(para)
        doc.save(path)

        text, page_count, is_scanned = extractor._extract_odt(path)

        print(f"\n--- Extracted ODT text ---\n{text}\n--- end ---")

        checks = []
        checks.append(("'# Introduction' in text", "# Introduction" in text))
        checks.append(("'This is the body text.' in text", "This is the body text." in text))
        checks.append(("page_count == 0", page_count == 0))
        checks.append(("is_scanned == False", is_scanned == False))

        failed = [c for c in checks if not c[1]]
        detail = "\n".join(f"  {'OK' if c[1] else 'FAIL'}: {c[0]}" for c in checks)
        report("ODT Extraction", len(failed) == 0, detail)


# =====================================================================
# Test 5: .doc (expects graceful failure without antiword)
# =====================================================================
def test_doc():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.doc")
        # Write a dummy file (not a real .doc, but we just need the extractor to attempt it)
        with open(path, "wb") as f:
            f.write(b"\xd0\xcf\x11\xe0" + b"\x00" * 100)  # OLE magic bytes

        # Import the private exception
        from etl.attachment_text_extractor import _UnsupportedType

        try:
            text, page_count, is_scanned = extractor._extract_doc(path)
            # If antiword IS installed, it might succeed or fail with RuntimeError
            # Either way, we got a result -- check it's a tuple
            report("DOC Extraction (antiword present)",
                   True,
                   f"  antiword is installed, got result: ({len(text)} chars, {page_count}, {is_scanned})")
        except _UnsupportedType as e:
            report("DOC Extraction (antiword not installed - graceful failure)",
                   "antiword not installed" in str(e),
                   f"  Correctly raised _UnsupportedType: {e}")
        except RuntimeError as e:
            # antiword is installed but failed on our dummy file -- that's expected
            report("DOC Extraction (antiword present, bad file)",
                   "antiword failed" in str(e),
                   f"  antiword present but rejected dummy file (expected): {e}")


# =====================================================================
# Run all tests
# =====================================================================
if __name__ == "__main__":
    tests = [
        ("PPTX", test_pptx),
        ("XLS", test_xls),
        ("RTF", test_rtf),
        ("ODT", test_odt),
        ("DOC", test_doc),
    ]

    for name, fn in tests:
        try:
            fn()
        except Exception as e:
            report(f"{name} Extraction (EXCEPTION)", False, traceback.format_exc())

    print(f"\n\n{'='*60}")
    print("SUMMARY")
    print(f"{'='*60}")
    passed = sum(1 for _, p in results if p)
    failed = sum(1 for _, p in results if not p)
    for name, p in results:
        print(f"  {'PASS' if p else 'FAIL'}: {name}")
    print(f"\nTotal: {passed} passed, {failed} failed out of {len(results)}")
    sys.exit(0 if failed == 0 else 1)
