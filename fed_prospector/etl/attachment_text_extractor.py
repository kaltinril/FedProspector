"""Structure-aware text extraction from opportunity attachments (Phase 110).

Downloads are already on disk (Phase 110 download step). This module extracts
text content from PDF, Word, Excel, and plain-text attachments, converting to
annotated Markdown with heading hierarchy, bold markers, and table formatting.

Usage:
    from etl.attachment_text_extractor import AttachmentTextExtractor
    extractor = AttachmentTextExtractor()
    stats = extractor.extract_text(batch_size=50)
"""

import hashlib
import logging
import os
import threading
from datetime import datetime
from pathlib import Path

from config.settings import ATTACHMENT_DIR as _DEFAULT_ATTACHMENT_DIR
from db.connection import get_connection
from etl.load_manager import LoadManager

logger = logging.getLogger("fed_prospector.etl.attachment_text_extractor")

# Suppress noisy MuPDF warnings about malformed PDF structure trees.
# These are non-fatal — text extraction succeeds despite them.
try:
    import pymupdf as fitz
    fitz.TOOLS.mupdf_display_errors(False)
except Exception:
    pass  # Older PyMuPDF versions may not have this

# Suppress openpyxl warnings about unsupported Excel features
# (data validation, print areas, etc.) — text extraction still works.
import warnings
warnings.filterwarnings("ignore", module="openpyxl")

# Content type / extension mappings
PDF_TYPES = {"application/pdf"}
WORD_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
EXCEL_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
}
TEXT_TYPES = {"text/plain", "text/csv", "text/html"}
PPTX_TYPES = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
RTF_TYPES = {"application/rtf", "text/rtf"}
ODT_TYPES = {"application/vnd.oasis.opendocument.text"}

PDF_EXTENSIONS = {".pdf"}
WORD_EXTENSIONS = {".docx", ".doc"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
PPTX_EXTENSIONS = {".pptx"}
RTF_EXTENSIONS = {".rtf"}
ODT_EXTENSIONS = {".odt"}
TEXT_EXTENSIONS = {".txt", ".csv", ".html", ".htm"}

# PyMuPDF font flag bitmask for bold
BOLD_FLAG = 16

# Minimum text chars per page before considering it scanned/image-only
SCANNED_THRESHOLD = 50


def _cell_str(value):
    """Convert a cell value to a clean string for Markdown tables."""
    if value is None:
        return ""
    return str(value).replace("|", "\\|").replace("\n", " ").strip()


# ------------------------------------------------------------------
# Module-level extraction handlers (picklable for ProcessPoolExecutor)
# ------------------------------------------------------------------

def _resolve_handler(content_type, ext):
    """Return the appropriate extraction function, or None if unsupported."""
    if content_type in PDF_TYPES or ext in PDF_EXTENSIONS:
        return _extract_pdf
    if content_type in PPTX_TYPES or ext in PPTX_EXTENSIONS:
        return _extract_pptx
    if content_type in WORD_TYPES or ext in WORD_EXTENSIONS:
        if ext == ".doc" or (content_type == "application/msword" and ext != ".docx"):
            return _extract_doc
        return _extract_docx
    if content_type in EXCEL_TYPES or ext in EXCEL_EXTENSIONS:
        if ext == ".xls" or (content_type == "application/vnd.ms-excel" and ext != ".xlsx"):
            return _extract_xls
        return _extract_xlsx
    if content_type in RTF_TYPES or ext in RTF_EXTENSIONS:
        return _extract_rtf
    if content_type in ODT_TYPES or ext in ODT_EXTENSIONS:
        return _extract_odt
    if content_type in TEXT_TYPES or ext in TEXT_EXTENSIONS:
        return _extract_plain_text
    return None


def _detect_type_by_magic(file_path):
    """Detect file type from magic bytes, ignoring extension/content_type.

    Returns handler function or None.
    """
    try:
        with open(file_path, "rb") as f:
            header = f.read(8)
    except Exception:
        return None

    if header.startswith(b"%PDF"):
        return _extract_pdf

    if header.startswith(b"PK\x03\x04"):
        # ZIP-based format -- try to distinguish by examining ZIP contents
        import zipfile
        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                names = zf.namelist()
                if any(n.startswith("word/") for n in names):
                    return _extract_docx
                if any(n.startswith("ppt/") for n in names):
                    return _extract_pptx
                if any(n.startswith("xl/") for n in names):
                    return _extract_xlsx
                if "content.xml" in names:
                    return _extract_odt
        except zipfile.BadZipFile:
            return None
        return None

    if header.startswith(b"\xd0\xcf\x11\xe0"):
        # OLE2 -- could be .doc or .xls. Try doc first (more common in solicitations)
        return _extract_doc

    if header.startswith(b"{\\rtf"):
        return _extract_rtf

    return None


def _check_ole2_encryption(file_path):
    """Raise early if an OLE2 file is IRM/DRM-encrypted.

    Checks for EncryptedPackage or DataSpaces streams that indicate
    Microsoft Information Rights Management protection.
    """
    try:
        with open(file_path, "rb") as f:
            header = f.read(8)
        if not header.startswith(b"\xd0\xcf\x11\xe0"):
            return  # Not OLE2, skip check

        import olefile
        if not olefile.isOleFile(file_path):
            return
        ole = olefile.OleFileIO(file_path)
        try:
            streams = {"/".join(s) for s in ole.listdir()}
            if "EncryptedPackage" in streams or "\x06DataSpaces" in streams:
                raise RuntimeError(
                    "File is IRM/DRM-protected (Microsoft Information Rights Management)"
                )
        finally:
            ole.close()
    except RuntimeError:
        raise
    except Exception:
        pass  # If olefile isn't installed or parsing fails, let LibreOffice try


# ------------------------------------------------------------------
# PDF extraction (PyMuPDF / fitz)
# ------------------------------------------------------------------

def _extract_pdf(file_path):
    """Extract text from PDF with structure-aware Markdown formatting.

    Returns:
        (markdown_text, page_count, is_scanned)
    """
    import pymupdf as fitz  # PyMuPDF 1.27+

    try:
        doc = fitz.open(file_path)
    except Exception as e:
        if "encrypted" in str(e).lower() or "password" in str(e).lower():
            raise RuntimeError(f"Password-protected PDF: {e}") from e
        raise

    page_count = len(doc)
    doc.close()

    # Table detection: safe to run now since each file is in its own process.
    # A hang or crash only kills that one worker, not the whole batch.
    if page_count <= 30:
        table_timeout = 10  # 10s per page
    else:
        table_timeout = 0  # skip for large PDFs (too many pages)

    return _extract_pdf_pages(file_path, range(page_count), table_timeout=table_timeout)


def _extract_pdf_pages(file_path, page_range, table_timeout=10):
    """Serial PDF extraction for small documents."""
    import pymupdf as fitz

    doc = fitz.open(file_path)
    pages = []
    is_scanned = False

    try:
        for page_idx in page_range:
            page = doc[page_idx]
            page_num = page_idx + 1

            table_rects = []
            table_md = ""
            if table_timeout > 0 and hasattr(page, "find_tables"):
                table_rects, table_md = _find_tables_with_timeout(file_path, page_idx, timeout=table_timeout)

            try:
                blocks = page.get_text("dict", sort=True)["blocks"]
            except Exception:
                blocks = []

            page_text = _process_pdf_blocks(blocks, table_rects)

            if table_md:
                page_text = page_text.rstrip() + "\n\n" + table_md

            if len(page_text.strip()) < SCANNED_THRESHOLD:
                is_scanned = True
                page_text = f"[Page {page_num}: scanned/image-only — OCR not available]\n"

            pages.append(page_text)
    finally:
        doc.close()

    return "\n\n---\n\n".join(pages), len(pages), is_scanned


def _find_tables_with_timeout(file_path, page_idx, timeout=10):
    """Run find_tables() with a try/except.

    Since extraction runs in a subprocess (ProcessPoolExecutor), a crash
    here only kills this worker, not the main process. No need for the
    old multiprocessing.Process wrapper.

    Returns:
        (table_rects, table_md) — or ([], "") on error.
    """
    try:
        import pymupdf as fitz
        doc = fitz.open(file_path)
        try:
            page = doc[page_idx]
            tables = page.find_tables()
            rects = []
            md_parts = []
            for table in tables:
                rects.append(tuple(table.bbox))
                try:
                    data = table.extract()
                except Exception:
                    continue
                if not data or not data[0]:
                    continue
                header = data[0]
                col_count = len(header)
                lines = []
                lines.append("| " + " | ".join(_cell_str(c) for c in header) + " |")
                lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")
                for row in data[1:]:
                    cells = list(row) + [""] * max(0, col_count - len(row))
                    lines.append("| " + " | ".join(_cell_str(c) for c in cells[:col_count]) + " |")
                md_parts.append("\n".join(lines))
            return rects, "\n\n".join(md_parts)
        finally:
            doc.close()
    except Exception:
        return [], ""


def _process_pdf_blocks(blocks, table_rects=None):
    """Convert PDF dict blocks to Markdown text.

    Uses font size to detect headings and font flags for bold.
    Skips blocks that overlap with already-extracted table regions.
    """
    table_rects = table_rects or []
    lines_out = []

    for block in blocks:
        # Skip image blocks
        if block.get("type") != 0:
            continue

        # Skip blocks inside table regions
        bx0, by0, bx1, by1 = block["bbox"]
        if _rect_overlaps_any(bx0, by0, bx1, by1, table_rects):
            continue

        for line in block.get("lines", []):
            spans = line.get("spans", [])
            if not spans:
                continue

            # Determine dominant font properties from first span
            first_span = spans[0]
            font_size = first_span.get("size", 11)
            is_bold = bool(first_span.get("flags", 0) & BOLD_FLAG)

            # Build line text from all spans
            line_text = ""
            for span in spans:
                text = span.get("text", "")
                if not text.strip():
                    line_text += text
                    continue
                span_bold = bool(span.get("flags", 0) & BOLD_FLAG)
                if span_bold and not (font_size >= 14 or (font_size >= 12 and is_bold)):
                    # Inline bold (not already a heading)
                    line_text += f"**{text.strip()}** " if text.strip() else text
                else:
                    line_text += text

            line_text = line_text.rstrip()
            if not line_text.strip():
                continue

            # Heading detection based on font size
            if font_size >= 14:
                lines_out.append(f"\n## {line_text.strip()}\n")
            elif font_size >= 12 and is_bold:
                lines_out.append(f"\n### {line_text.strip()}\n")
            else:
                lines_out.append(line_text)

    return "\n".join(lines_out)


def _pdf_table_to_markdown(table):
    """Convert a PyMuPDF table object to Markdown table format."""
    try:
        data = table.extract()
    except Exception:
        return ""

    if not data or not data[0]:
        return ""

    # First row as header
    header = data[0]
    col_count = len(header)
    lines = []
    lines.append("| " + " | ".join(_cell_str(c) for c in header) + " |")
    lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")

    for row in data[1:]:
        # Pad or trim to match header column count
        cells = list(row) + [""] * max(0, col_count - len(row))
        lines.append("| " + " | ".join(_cell_str(c) for c in cells[:col_count]) + " |")

    return "\n".join(lines)


def _rect_overlaps_any(x0, y0, x1, y1, rects):
    """Check if a bounding box overlaps with any rect in the list."""
    for rx0, ry0, rx1, ry1 in rects:
        if x0 < rx1 and x1 > rx0 and y0 < ry1 and y1 > ry0:
            return True
    return False


# ------------------------------------------------------------------
# Word extraction (python-docx)
# ------------------------------------------------------------------

def _extract_docx(file_path):
    """Extract text from .docx with heading and bold detection.

    Returns:
        (markdown_text, page_count=0, is_scanned=False)
    """
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(file_path)
    parts = []

    # Build O(1) lookup maps instead of scanning all paragraphs/tables per element
    para_map = {id(para._element): para for para in doc.paragraphs}
    table_map = {id(tbl._element): tbl for tbl in doc.tables}

    # Iterate body children to maintain paragraph/table ordering
    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            para = para_map.get(id(element))
            if para is not None:
                md = _docx_paragraph_to_md(para)
                if md:
                    parts.append(md)

        elif tag == "tbl":
            table = table_map.get(id(element))
            if table is not None:
                md = _docx_table_to_md(table)
                if md:
                    parts.append(md)

    text = "\n".join(parts)
    # python-docx has no page count concept
    return text, 0, False


def _docx_paragraph_to_md(para):
    """Convert a python-docx Paragraph to Markdown."""
    style_name = (para.style.name or "").lower() if para.style else ""

    # Build text with bold markers
    text_parts = []
    for run in para.runs:
        t = run.text
        if not t:
            continue
        if run.bold:
            text_parts.append(f"**{t.strip()}**" if t.strip() else t)
        else:
            text_parts.append(t)

    text = "".join(text_parts).strip()
    if not text:
        return ""

    # Heading detection from style
    if "heading 1" in style_name:
        return f"\n# {text}\n"
    if "heading 2" in style_name:
        return f"\n## {text}\n"
    if "heading 3" in style_name or "heading 4" in style_name:
        return f"\n### {text}\n"
    if "title" in style_name:
        return f"\n# {text}\n"

    return text


def _docx_table_to_md(table):
    """Convert a python-docx Table to Markdown."""
    rows = table.rows
    if not rows:
        return ""

    md_rows = []
    col_count = len(rows[0].cells)

    for i, row in enumerate(rows):
        cells = [_cell_str(cell.text) for cell in row.cells[:col_count]]
        md_rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            md_rows.append("| " + " | ".join("---" for _ in range(col_count)) + " |")

    return "\n".join(md_rows)


# ------------------------------------------------------------------
# Excel extraction (openpyxl)
# ------------------------------------------------------------------

def _extract_xlsx(file_path):
    """Extract text from .xlsx with sheet and table structure.

    Returns:
        (markdown_text, page_count=0, is_scanned=False)
    """
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))

        # Skip empty sheets
        if not rows or all(all(c is None for c in row) for row in rows):
            continue

        parts.append(f"\n## {sheet_name}\n")

        # Filter out completely empty rows
        data_rows = [row for row in rows if any(c is not None for c in row)]
        if not data_rows:
            continue

        # Determine column count from first row
        col_count = len(data_rows[0])

        # First row as header
        header = data_rows[0]
        parts.append("| " + " | ".join(_cell_str(c) for c in header) + " |")
        parts.append("| " + " | ".join("---" for _ in range(col_count)) + " |")

        for row in data_rows[1:]:
            cells = list(row) + [None] * max(0, col_count - len(row))
            parts.append("| " + " | ".join(_cell_str(c) for c in cells[:col_count]) + " |")

    wb.close()
    return "\n".join(parts), 0, False


# ------------------------------------------------------------------
# Legacy Word extraction (.doc -> .docx via LibreOffice)
# ------------------------------------------------------------------

def _extract_doc(file_path):
    """Extract text from legacy .doc by converting to .docx via LibreOffice.

    Converts the .doc to a temporary .docx, then delegates to _extract_docx
    for full structure-aware extraction (headings, bold, tables).

    Returns:
        (markdown_text, page_count=0, is_scanned=False)
    """
    import shutil
    import subprocess
    import tempfile

    soffice = shutil.which("soffice")
    if not soffice:
        # Check common install path on Windows
        win_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        if Path(win_path).is_file():
            soffice = win_path

    if not soffice:
        logger.warning(
            "LibreOffice not installed — .doc extraction unavailable. "
            "Install LibreOffice to enable legacy Word support."
        )
        raise _UnsupportedType("LibreOffice not installed for .doc extraction")

    # Detect IRM/DRM-protected files before wasting time on LibreOffice
    _check_ole2_encryption(file_path)

    with tempfile.TemporaryDirectory() as tmpdir:
        # LibreOffice refuses to convert OLE2 files with a .docx extension.
        # Copy to a .doc temp file to ensure correct handling.
        input_path = file_path
        if not file_path.lower().endswith(".doc"):
            tmp_input = os.path.join(tmpdir, "input.doc")
            shutil.copy2(file_path, tmp_input)
            input_path = tmp_input

        # Clean environment to prevent Python venv from interfering
        # with LibreOffice's bundled Python
        env = {k: v for k, v in os.environ.items()
               if k not in ("PYTHONPATH", "PYTHONHOME", "VIRTUAL_ENV")}

        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "docx",
             "--outdir", tmpdir, input_path],
            capture_output=True, text=True, timeout=60, env=env,
        )
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr.strip()}")

        # Find the converted .docx file
        docx_files = list(Path(tmpdir).glob("*.docx"))
        if not docx_files:
            # Check for IRM/DRM protected files
            stderr = result.stderr or ""
            if "source file could not be loaded" in stderr.lower():
                raise RuntimeError("File could not be loaded (possibly IRM/DRM protected)")
            raise RuntimeError("LibreOffice produced no .docx output")

        return _extract_docx(str(docx_files[0]))


# ------------------------------------------------------------------
# PowerPoint extraction (python-pptx)
# ------------------------------------------------------------------

def _extract_pptx(file_path):
    """Extract text from .pptx with slide headings, tables, and speaker notes.

    Returns:
        (markdown_text, page_count=slide_count, is_scanned=False)
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = [f"\n## Slide {slide_num}\n"]

        for shape in slide.shapes:
            if shape.has_table:
                slide_parts.append(_pptx_table_to_md(shape.table))
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"\n> **Notes:** {notes_text}\n")

        parts.append("\n".join(slide_parts))

    return "\n\n".join(parts), len(prs.slides), False


def _pptx_table_to_md(table):
    """Convert a python-pptx Table to Markdown."""
    rows = list(table.rows)
    if not rows:
        return ""

    col_count = len(table.columns)
    md_rows = []

    for i, row in enumerate(rows):
        cells = [_cell_str(cell.text) for cell in list(row.cells)[:col_count]]
        md_rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            md_rows.append("| " + " | ".join("---" for _ in range(col_count)) + " |")

    return "\n".join(md_rows)


# ------------------------------------------------------------------
# Legacy Excel extraction (xlrd)
# ------------------------------------------------------------------

def _extract_xls(file_path):
    """Extract text from legacy .xls files using xlrd.

    Returns:
        (markdown_text, page_count=0, is_scanned=False)
    """
    import xlrd

    wb = xlrd.open_workbook(file_path)
    parts = []

    for sheet in wb.sheets():
        if sheet.nrows == 0:
            continue

        # Skip sheets where all cells are empty
        has_data = False
        for row_idx in range(sheet.nrows):
            if any(sheet.cell_value(row_idx, col_idx) for col_idx in range(sheet.ncols)):
                has_data = True
                break
        if not has_data:
            continue

        parts.append(f"\n## {sheet.name}\n")

        # Collect non-empty rows
        data_rows = []
        for row_idx in range(sheet.nrows):
            row = [sheet.cell_value(row_idx, col_idx) for col_idx in range(sheet.ncols)]
            if any(c is not None and c != "" for c in row):
                data_rows.append(row)

        if not data_rows:
            continue

        col_count = len(data_rows[0])

        # First row as header
        header = data_rows[0]
        parts.append("| " + " | ".join(_cell_str(c) for c in header) + " |")
        parts.append("| " + " | ".join("---" for _ in range(col_count)) + " |")

        for row in data_rows[1:]:
            cells = list(row) + [""] * max(0, col_count - len(row))
            parts.append("| " + " | ".join(_cell_str(c) for c in cells[:col_count]) + " |")

    return "\n".join(parts), 0, False


# ------------------------------------------------------------------
# RTF extraction (striprtf)
# ------------------------------------------------------------------

def _extract_rtf(file_path):
    """Extract text from .rtf files using striprtf.

    Returns:
        (text, page_count=0, is_scanned=False)
    """
    from striprtf.striprtf import rtf_to_text

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        rtf_content = f.read()

    text = rtf_to_text(rtf_content)
    return text, 0, False


# ------------------------------------------------------------------
# ODT extraction (odfpy)
# ------------------------------------------------------------------

def _extract_odt(file_path):
    """Extract text from .odt files using odfpy.

    Returns:
        (markdown_text, page_count=0, is_scanned=False)
    """
    from odf.opendocument import load
    from odf.text import H, P
    from odf.table import Table, TableRow, TableCell
    from odf import teletype

    doc = load(file_path)
    parts = []

    for element in doc.text.childNodes:
        tag = element.qname[1] if hasattr(element, 'qname') else ""

        if tag == "h":
            # Heading — detect outline level
            level = element.getAttribute("outlinelevel") or "2"
            try:
                level = int(level)
            except (ValueError, TypeError):
                level = 2
            prefix = "#" * min(level, 4)
            text = teletype.extractText(element).strip()
            if text:
                parts.append(f"\n{prefix} {text}\n")

        elif tag == "p":
            text = teletype.extractText(element).strip()
            if text:
                parts.append(text)

        elif tag == "table":
            md = _odt_table_to_md(element)
            if md:
                parts.append(md)

    return "\n".join(parts), 0, False


def _odt_table_to_md(table_element):
    """Convert an ODF table element to Markdown."""
    from odf.table import TableRow, TableCell
    from odf import teletype

    rows = []
    for row_el in table_element.getElementsByType(TableRow):
        cells = []
        for cell_el in row_el.getElementsByType(TableCell):
            cells.append(_cell_str(teletype.extractText(cell_el)))
        if cells:
            rows.append(cells)

    if not rows:
        return ""

    col_count = len(rows[0])
    md_rows = []

    for i, row in enumerate(rows):
        cells = list(row) + [""] * max(0, col_count - len(row))
        md_rows.append("| " + " | ".join(cells[:col_count]) + " |")
        if i == 0:
            md_rows.append("| " + " | ".join("---" for _ in range(col_count)) + " |")

    return "\n".join(md_rows)


# ------------------------------------------------------------------
# Plain text extraction
# ------------------------------------------------------------------

def _extract_plain_text(file_path):
    """Read plain text files directly.

    Returns:
        (text, page_count=0, is_scanned=False)
    """
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                text = f.read()
            return text, 0, False
        except (UnicodeDecodeError, UnicodeError):
            continue

    # Last resort: read as bytes and decode lossy
    with open(file_path, "rb") as f:
        raw = f.read()
    return raw.decode("utf-8", errors="replace"), 0, False


# ------------------------------------------------------------------
# Subprocess entry point for ProcessPoolExecutor
# ------------------------------------------------------------------

def _extract_file(file_path, content_type, ext, filename, attachment_dir):
    """Extract text from a single file. Runs in a subprocess.

    Returns a dict with status and extraction results:
        {"status": "extracted", "text": ..., "page_count": ..., "is_scanned": ...}
        {"status": "unsupported"}
        {"status": "failed", "error": str}
    """
    try:
        # Resolve full path
        full_path = Path(attachment_dir) / file_path
        if not full_path.is_file():
            return {"status": "failed", "error": f"File not found: {full_path}"}
        file_path_str = str(full_path)

        handler = _resolve_handler(content_type, ext)

        if handler is None:
            # Try magic byte detection as fallback
            handler = _detect_type_by_magic(file_path_str)

        if handler is None:
            return {"status": "unsupported"}

        try:
            text, page_count, is_scanned = handler(file_path_str)
        except _UnsupportedType:
            return {"status": "unsupported"}
        except Exception as first_err:
            # If handler failed, try magic byte detection in case extension was wrong
            magic_handler = _detect_type_by_magic(file_path_str)
            if magic_handler and magic_handler != handler:
                text, page_count, is_scanned = magic_handler(file_path_str)
            else:
                raise

        if not text or not text.strip():
            text = "[No extractable text content]"
            # Only preserve is_scanned when the handler actually inspected pages
            # for image-only content (PDF). For other formats, an empty extraction
            # means the file genuinely has no text — not that it's scanned.
            if handler is not _extract_pdf:
                is_scanned = False

        text_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()

        return {
            "status": "extracted",
            "text": text,
            "text_hash": text_hash,
            "page_count": page_count,
            "is_scanned": is_scanned,
        }

    except _UnsupportedType:
        return {"status": "unsupported"}
    except Exception as e:
        return {"status": "failed", "error": str(e)}


# ------------------------------------------------------------------
# Main class
# ------------------------------------------------------------------

class AttachmentTextExtractor:
    """Extracts text from downloaded opportunity attachments.

    Produces annotated Markdown with heading hierarchy, bold markers,
    and table formatting. Stores results back in the attachment_document
    table.
    """

    def __init__(self, db_connection=None, load_manager=None, attachment_dir=None):
        self.db_connection = db_connection
        self.load_manager = load_manager or LoadManager()
        self.attachment_dir = Path(attachment_dir) if attachment_dir else _DEFAULT_ATTACHMENT_DIR

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def extract_text(self, notice_id=None, batch_size=100, force=False, workers=10,
                     timeout_seconds=120):
        """Extract text from pending downloaded attachments.

        Args:
            notice_id: If set, only process attachments for this notice.
            batch_size: Max attachments to process per run.
            force: If True, re-extract even if already extracted.
            workers: Number of concurrent file extraction processes (default 10).
            timeout_seconds: Per-file wall-clock timeout in seconds (default 120).
                When exceeded, the worker is SIGTERM'd, the document is marked
                extraction_status='timeout', and the pool continues with a
                respawned worker.

        Returns:
            dict with keys: processed, extracted, failed, unsupported, timeout, skipped
        """
        load_id = self.load_manager.start_load(
            source_system="ATTACHMENT_TEXT",
            load_type="INCREMENTAL",
            parameters={
                "notice_id": notice_id,
                "batch_size": batch_size,
                "force": force,
                "timeout_seconds": timeout_seconds,
            },
        )

        stats = {
            "processed": 0,
            "extracted": 0,
            "failed": 0,
            "unsupported": 0,
            "timeout": 0,
            "skipped": 0,
        }

        # Phase 124: per-run dedup state — reset every run, not in __init__.
        # _seen_text_hashes maps text_hash -> canonical attachment_id seen in
        # this batch (defends against duplicates within a single run).
        self._seen_text_hashes = {}
        self._seen_lock = threading.Lock()
        self._layer4_dedup_count = 0
        self._in_batch_dedup_count = 0

        try:
            rows = self._fetch_pending(notice_id, batch_size, force)
            logger.info("Found %d attachments to extract", len(rows))

            if not rows:
                self.load_manager.complete_load(
                    load_id, records_read=0, records_inserted=0,
                    records_updated=0, records_unchanged=0, records_errored=0,
                )
                return {"processed": 0, "succeeded": 0, "failed": 0, "skipped": 0}

            from concurrent.futures import TimeoutError as FuturesTimeoutError, as_completed
            from pebble import ProcessPool
            from tqdm import tqdm

            pbar = tqdm(
                total=len(rows),
                desc="Extracting",
                unit="file",
                bar_format="{desc}: {n_fmt}/{total_fmt} [{elapsed}<{remaining}] {postfix}",
            )

            actual_workers = min(workers, len(rows))
            attachment_dir_str = str(self.attachment_dir)

            # Use pebble.ProcessPool (not concurrent.futures.ProcessPoolExecutor)
            # because it supports per-task timeouts that actually SIGTERM stuck
            # workers and respawn them. stdlib ProcessPoolExecutor cannot cancel
            # running tasks and future.result(timeout=N) inside as_completed is
            # a no-op.
            with ProcessPool(max_workers=actual_workers) as pool:
                futures = {}
                for row in rows:
                    file_path = row.get("file_path")
                    filename = row.get("filename") or ""
                    content_type = (row.get("content_type") or "").lower()
                    ext = os.path.splitext(filename)[1].lower() if filename else ""

                    future = pool.schedule(
                        _extract_file,
                        args=(file_path, content_type, ext, filename, attachment_dir_str),
                        timeout=timeout_seconds,
                    )
                    futures[future] = row

                # pebble.ProcessFuture is compatible with concurrent.futures.as_completed
                for future in as_completed(list(futures.keys())):
                    row = futures[future]
                    notice_id_val = row.get("notice_id") or "?"
                    filename = row.get("filename") or "unknown"
                    attachment_id = row["attachment_id"]

                    timed_out = False
                    try:
                        result = future.result()
                    except FuturesTimeoutError:
                        # Worker was SIGTERM'd by pebble after timeout_seconds.
                        # Pool has already respawned the replacement worker.
                        timed_out = True
                        result = None
                    except Exception as e:
                        # Subprocess crashed hard (should be rare)
                        result = {"status": "failed", "error": str(e)}

                    if timed_out:
                        logger.warning(
                            "Extraction timeout after %ds for attachment %s (%s, notice=%s)",
                            timeout_seconds, attachment_id, filename, notice_id_val,
                        )
                        try:
                            self._mark_timeout(attachment_id, load_id)
                        except Exception as mark_err:
                            logger.error(
                                "Failed to mark timeout for attachment %s: %s",
                                attachment_id, mark_err,
                            )
                        self.load_manager.log_record_error(
                            load_id,
                            record_identifier=str(attachment_id),
                            error_type="EXTRACTION_TIMEOUT",
                            error_message=f"Timed out after {timeout_seconds}s",
                        )
                        stats["timeout"] += 1
                    elif result["status"] == "extracted":
                        try:
                            self._save_extraction(
                                attachment_id=attachment_id,
                                text=result["text"],
                                text_hash=result["text_hash"],
                                page_count=result["page_count"],
                                is_scanned=result["is_scanned"],
                                load_id=load_id,
                            )
                            stats["extracted"] += 1
                            # Phase 124 Layer 4: text-hash dedup check after
                            # successful extraction. If this document is a
                            # duplicate of an already-processed canonical, the
                            # extracted row, file, and opportunity mappings are
                            # remapped to the canonical and this document is
                            # removed before downstream intel sees it.
                            try:
                                self._handle_text_hash_dedup(
                                    document_id=attachment_id,
                                    text_hash=result["text_hash"],
                                )
                            except Exception as dedup_err:
                                logger.error(
                                    "Layer 4 text-hash dedup failed for document %s: %s",
                                    attachment_id, dedup_err,
                                )
                        except Exception as e:
                            self._mark_failed(attachment_id, load_id, str(e))
                            stats["failed"] += 1
                    elif result["status"] == "unsupported":
                        self._mark_unsupported(attachment_id, load_id)
                        stats["unsupported"] += 1
                    else:
                        # failed
                        error_msg = result.get("error", "Unknown error")
                        self._mark_failed(attachment_id, load_id, error_msg)
                        self.load_manager.log_record_error(
                            load_id,
                            record_identifier=str(attachment_id),
                            error_type="EXTRACTION_ERROR",
                            error_message=error_msg,
                        )
                        logger.error(
                            "Extraction failed for attachment %s (%s): %s",
                            attachment_id, filename, error_msg,
                        )
                        stats["failed"] += 1

                    stats["processed"] += 1
                    pbar.set_postfix_str(
                        f"ok={stats['extracted']} fail={stats['failed']} unsup={stats['unsupported']} timeout={stats['timeout']} | {notice_id_val} | {filename}"
                    )
                    pbar.update(1)

            pbar.close()

            self.load_manager.complete_load(
                load_id,
                records_read=stats["processed"],
                records_inserted=stats["extracted"],
                records_updated=0,
                records_unchanged=stats["skipped"],
                records_errored=stats["failed"] + stats["unsupported"] + stats["timeout"],
            )
            logger.info(
                "Extraction complete: %d extracted, %d failed, %d unsupported, %d timeout, %d skipped",
                stats["extracted"],
                stats["failed"],
                stats["unsupported"],
                stats["timeout"],
                stats["skipped"],
            )
            # Phase 124 Task 9: dedup summary for the run.
            logger.info(
                "Dedup: layer4_text_hash_skipped=%d, in_batch_skipped=%d",
                self._layer4_dedup_count,
                self._in_batch_dedup_count,
            )
            stats["layer4_text_hash_skipped"] = self._layer4_dedup_count
            stats["in_batch_skipped"] = self._in_batch_dedup_count
        except Exception as e:
            self.load_manager.fail_load(load_id, str(e))
            logger.error("Extraction batch failed: %s", e)
            raise

        return stats

    # ------------------------------------------------------------------
    # Single-attachment re-extraction
    # ------------------------------------------------------------------

    def reextract_single(self, attachment_id):
        """Re-extract text from a single attachment by sam_attachment.attachment_id.

        Looks up the attachment's file path, re-runs text extraction, and
        updates the attachment_document row.

        Returns:
            dict with keys: processed, extracted, failed
        """
        conn = get_connection()
        cursor = conn.cursor(dictionary=True)
        try:
            cursor.execute(
                "SELECT ad.document_id, ad.attachment_id, "
                "COALESCE(ad.filename, sa.filename) AS filename, "
                "ad.content_type, sa.file_path "
                "FROM attachment_document ad "
                "JOIN sam_attachment sa ON sa.attachment_id = ad.attachment_id "
                "WHERE sa.attachment_id = %s",
                (attachment_id,),
            )
            row = cursor.fetchone()
        finally:
            cursor.close()
            conn.close()

        if not row:
            raise ValueError(f"No attachment_document found for attachment_id {attachment_id}")

        load_id = self.load_manager.start_load(
            source_system="ATTACHMENT_TEXT",
            load_type="INCREMENTAL",
            parameters={"attachment_id": attachment_id, "reextract": True},
        )

        stats = {"processed": 1, "extracted": 0, "failed": 0}

        try:
            # Use document_id as attachment_id (same convention as _fetch_pending)
            row["attachment_id"] = row["document_id"]
            self._extract_one(row, load_id)
            stats["extracted"] = 1
        except Exception as e:
            stats["failed"] = 1
            logger.error("Re-extraction failed for attachment %d: %s", attachment_id, e)

        self.load_manager.complete_load(
            load_id,
            records_read=1,
            records_inserted=stats["extracted"],
            records_errored=stats["failed"],
        )
        return stats

    # ------------------------------------------------------------------
    # Internal: query and dispatch
    # ------------------------------------------------------------------

    def _fetch_pending(self, notice_id, batch_size, force):
        """Fetch documents eligible for extraction.

        Queries attachment_document joined to sam_attachment for file_path.
        When filtering by --notice-id, joins through opportunity_attachment map.
        """
        conn = self.db_connection or get_connection()
        cursor = conn.cursor(dictionary=True)
        try:
            conditions = ["sa.download_status = 'downloaded'"]
            params = []

            if force:
                conditions.append("ad.extraction_status IN ('pending', 'extracted', 'failed', 'unsupported')")
            else:
                conditions.append("ad.extraction_status = 'pending'")
                conditions.append("ad.extraction_retry_count < 10")

            join_clause = ""
            if notice_id:
                join_clause = (
                    "JOIN opportunity_attachment m ON m.attachment_id = sa.attachment_id "
                )
                conditions.append("m.notice_id = %s")
                params.append(notice_id)

            sql = (
                "SELECT ad.document_id, ad.attachment_id, "
                "COALESCE(ad.filename, sa.filename) AS filename, "
                "ad.content_type, sa.file_path "
                "FROM attachment_document ad "
                "JOIN sam_attachment sa ON sa.attachment_id = ad.attachment_id "
                + join_clause
                + f"WHERE {' AND '.join(conditions)} "
                "ORDER BY ad.document_id "
                "LIMIT %s"
            )
            params.append(batch_size)
            cursor.execute(sql, params)
            rows = cursor.fetchall()
            # Map document_id to attachment_id for downstream compatibility
            for row in rows:
                row["attachment_id"] = row["document_id"]
            return rows
        finally:
            cursor.close()
            if not self.db_connection:
                conn.close()

    def _extract_one(self, row, load_id):
        """Extract text from a single attachment and update the DB.

        Used for single-file extraction outside the process pool.
        """
        attachment_id = row["attachment_id"]
        file_path = row.get("file_path")
        filename = row.get("filename") or ""
        content_type = (row.get("content_type") or "").lower()
        ext = os.path.splitext(filename)[1].lower() if filename else ""

        result = _extract_file(file_path, content_type, ext, filename, str(self.attachment_dir))

        if result["status"] == "extracted":
            self._save_extraction(
                attachment_id=attachment_id,
                text=result["text"],
                text_hash=result["text_hash"],
                page_count=result["page_count"],
                is_scanned=result["is_scanned"],
                load_id=load_id,
            )
        elif result["status"] == "unsupported":
            self._mark_unsupported(attachment_id, load_id)
            raise _UnsupportedType(f"Unsupported: content_type={content_type}, ext={ext}")
        else:
            error_msg = result.get("error", "Unknown error")
            self._mark_failed(attachment_id, load_id, error_msg)
            raise RuntimeError(error_msg)

    # ------------------------------------------------------------------
    # DB persistence
    # ------------------------------------------------------------------

    def _save_extraction(self, attachment_id, text, text_hash, page_count, is_scanned, load_id):
        """Write extraction results to attachment_document.

        Note: attachment_id here is actually document_id (mapped in _fetch_pending).
        """
        document_id = attachment_id  # renamed for clarity
        conn = self.db_connection or get_connection()
        cursor = conn.cursor()
        try:
            # Strip null bytes — MySQL rejects them in TEXT columns
            clean_text = text.replace("\x00", "") if text else text
            cursor.execute(
                "UPDATE attachment_document SET "
                "extracted_text = %s, text_hash = %s, page_count = %s, "
                "is_scanned = %s, extraction_status = 'extracted', "
                "extracted_at = %s, last_load_id = %s "
                "WHERE document_id = %s",
                (
                    clean_text,
                    text_hash,
                    page_count if page_count else None,
                    1 if is_scanned else 0,
                    datetime.now(),
                    load_id,
                    document_id,
                ),
            )
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            cursor.close()
            if not self.db_connection:
                conn.close()

    # ------------------------------------------------------------------
    # Phase 124 Layer 4: text-hash dedup
    # ------------------------------------------------------------------

    def _handle_text_hash_dedup(self, document_id, text_hash):
        """Layer 4 text-hash dedup, called after successful extraction.

        If another already-processed `attachment_document` has the same
        `text_hash` (the canonical), this method:
          1. Remaps every `opportunity_attachment` row pointing at the current
             `sam_attachment.attachment_id` to the canonical `attachment_id`
             (INSERT IGNORE before DELETE so we never lose a reference).
          2. Records the mapping in `attachment_dedup_map` with method
             `'text_hash'` and both content_hash and text_hash populated.
          3. Deletes the just-extracted `attachment_document` row.
          4. Deletes the physical file from disk and clears
             `sam_attachment.file_path` so downstream cleanup leaves it alone.

        Also defends against in-batch duplicates: if two documents in the same
        run produce the same text_hash, the first wins, the second is treated
        the same way (the first becomes the canonical for the rest of the
        batch). All shared state is guarded by `self._seen_lock`.

        Filename-based intel is unaffected — `sam_attachment.filename` is
        preserved on the deduped row, so filename-only keyword matches still
        run on the canonical's downstream pass.

        Args:
            document_id: PK of the attachment_document row we just extracted.
            text_hash: SHA-256 of the extracted text.

        Returns:
            True if a duplicate was detected and handled, False otherwise.
        """
        if not text_hash:
            return False

        # Allow callers outside extract_text() (e.g. reextract_single) to use
        # this method without first initializing the per-run dedup state.
        if not hasattr(self, "_seen_lock") or self._seen_lock is None:
            self._seen_lock = threading.Lock()
        if not hasattr(self, "_seen_text_hashes") or self._seen_text_hashes is None:
            self._seen_text_hashes = {}
        if not hasattr(self, "_layer4_dedup_count"):
            self._layer4_dedup_count = 0
        if not hasattr(self, "_in_batch_dedup_count"):
            self._in_batch_dedup_count = 0

        conn = self.db_connection or get_connection()
        cursor = conn.cursor(dictionary=True)
        try:
            # Resolve the current document's sam_attachment_id and resource_guid.
            # We need both to remap opportunity_attachment and to record the
            # dedup map entry.
            cursor.execute(
                "SELECT ad.attachment_id AS sam_attachment_id, "
                "       sa.resource_guid AS resource_guid, "
                "       sa.content_hash AS content_hash, "
                "       sa.file_path AS file_path "
                "FROM attachment_document ad "
                "JOIN sam_attachment sa ON sa.attachment_id = ad.attachment_id "
                "WHERE ad.document_id = %s",
                (document_id,),
            )
            current = cursor.fetchone()
            if not current:
                return False

            current_sam_id = current["sam_attachment_id"]
            current_resource_guid = current["resource_guid"]
            current_content_hash = current["content_hash"]
            current_file_path = current["file_path"]

            # Find a canonical: a different attachment_document with the same
            # text_hash whose underlying attachment has been fully processed
            # (has at least one document_intel_summary row — same signal used
            # by attachment_cleanup to determine "intel complete").
            #
            # Also accept an in-batch canonical: the first document in this run
            # to record this text_hash wins, so subsequent ones in the same
            # batch dedup to it even before intel has run.
            canonical_sam_id = None
            with self._seen_lock:
                in_batch_canonical = self._seen_text_hashes.get(text_hash)
            if in_batch_canonical and in_batch_canonical != current_sam_id:
                canonical_sam_id = in_batch_canonical
                in_batch_hit = True
            else:
                in_batch_hit = False
                cursor.execute(
                    "SELECT ad.attachment_id "
                    "FROM attachment_document ad "
                    "WHERE ad.text_hash = %s "
                    "  AND ad.document_id <> %s "
                    "  AND ad.extraction_status = 'extracted' "
                    "  AND EXISTS ( "
                    "    SELECT 1 FROM document_intel_summary dis "
                    "    WHERE dis.document_id = ad.document_id "
                    "  ) "
                    "ORDER BY ad.document_id "
                    "LIMIT 1",
                    (text_hash, document_id),
                )
                match = cursor.fetchone()
                if match:
                    canonical_sam_id = match["attachment_id"]

            if not canonical_sam_id:
                # No duplicate found — record this document as the in-batch
                # canonical for any later siblings sharing the same text and
                # let the normal pipeline run on it.
                with self._seen_lock:
                    # Don't overwrite an existing entry — first one wins.
                    self._seen_text_hashes.setdefault(text_hash, current_sam_id)
                return False

            # Resolve canonical resource_guid (best-effort, audit-only).
            cursor.execute(
                "SELECT resource_guid FROM sam_attachment WHERE attachment_id = %s",
                (canonical_sam_id,),
            )
            canon_row = cursor.fetchone()
            canonical_resource_guid = canon_row["resource_guid"] if canon_row else None

            # Perform the remap + dedup writes in a single transaction so we
            # don't leave the DB in a half-deduped state on failure.
            write_cursor = conn.cursor()
            try:
                # 1. Remap opportunity_attachment: insert canonical row, then
                #    delete original. INSERT-before-DELETE preserves the link.
                write_cursor.execute(
                    "SELECT notice_id, url FROM opportunity_attachment "
                    "WHERE attachment_id = %s",
                    (current_sam_id,),
                )
                mappings = write_cursor.fetchall()
                for notice_id, url in mappings:
                    write_cursor.execute(
                        "INSERT IGNORE INTO opportunity_attachment "
                        "(notice_id, attachment_id, url) VALUES (%s, %s, %s)",
                        (notice_id, canonical_sam_id, url),
                    )
                    write_cursor.execute(
                        "DELETE FROM opportunity_attachment "
                        "WHERE notice_id = %s AND attachment_id = %s",
                        (notice_id, current_sam_id),
                    )

                # 2. Record the dedup decision so Layer 2 catches it on the
                #    next run without re-downloading.
                if current_resource_guid:
                    write_cursor.execute(
                        "INSERT IGNORE INTO attachment_dedup_map "
                        "(resource_guid, canonical_attachment_id, dedup_method, "
                        " content_hash, text_hash) "
                        "VALUES (%s, %s, 'text_hash', %s, %s)",
                        (
                            current_resource_guid,
                            canonical_sam_id,
                            current_content_hash,
                            text_hash,
                        ),
                    )

                # 3. Delete the just-extracted attachment_document row.
                write_cursor.execute(
                    "DELETE FROM attachment_document WHERE document_id = %s",
                    (document_id,),
                )

                # 4. Clear file_path on sam_attachment so cleanup leaves the
                #    deduped file alone going forward. The sam_attachment row
                #    stays as Layer 1's "I've seen this URL" record.
                write_cursor.execute(
                    "UPDATE sam_attachment SET file_path = NULL "
                    "WHERE attachment_id = %s",
                    (current_sam_id,),
                )

                conn.commit()
            except Exception:
                conn.rollback()
                raise
            finally:
                write_cursor.close()

            # 5. Best-effort: delete the physical file from disk. DB state is
            #    already correct even if this fails.
            if current_file_path:
                try:
                    full_path = self.attachment_dir / current_file_path
                    if full_path.is_file():
                        full_path.unlink()
                        # Try to remove empty parent dir (resource_guid folder)
                        try:
                            parent = full_path.parent
                            if parent != self.attachment_dir and not any(parent.iterdir()):
                                parent.rmdir()
                        except OSError:
                            pass
                except Exception as file_err:
                    logger.warning(
                        "Layer 4 dedup: failed to delete file %s for resource_guid=%s: %s",
                        current_file_path, current_resource_guid, file_err,
                    )

            # 6. Track + log.
            with self._seen_lock:
                if in_batch_hit:
                    self._in_batch_dedup_count += 1
                else:
                    self._layer4_dedup_count += 1
                # Future siblings should also dedup to the canonical, not to
                # this row (which we just deleted).
                self._seen_text_hashes[text_hash] = canonical_sam_id

            logger.info(
                "text-hash dedup: resource_guid=%s is duplicate of canonical "
                "attachment_id=%s (canonical_resource_guid=%s, method=text_hash, "
                "in_batch=%s)",
                current_resource_guid,
                canonical_sam_id,
                canonical_resource_guid,
                in_batch_hit,
            )
            return True
        finally:
            cursor.close()
            if not self.db_connection:
                conn.close()

    def _mark_failed(self, attachment_id, load_id, error_msg):
        """Mark document extraction as failed.

        Note: attachment_id here is actually document_id (mapped in _fetch_pending).
        """
        document_id = attachment_id
        conn = self.db_connection or get_connection()
        cursor = conn.cursor()
        try:
            cursor.execute(
                "UPDATE attachment_document SET "
                "extraction_status = 'failed', extraction_retry_count = extraction_retry_count + 1, "
                "extracted_at = %s, last_load_id = %s "
                "WHERE document_id = %s",
                (datetime.now(), load_id, document_id),
            )
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            cursor.close()
            if not self.db_connection:
                conn.close()

    def _mark_unsupported(self, attachment_id, load_id):
        """Mark document content type as unsupported.

        Note: attachment_id here is actually document_id (mapped in _fetch_pending).
        """
        document_id = attachment_id
        conn = self.db_connection or get_connection()
        cursor = conn.cursor()
        try:
            cursor.execute(
                "UPDATE attachment_document SET "
                "extraction_status = 'unsupported', extracted_at = %s, last_load_id = %s "
                "WHERE document_id = %s",
                (datetime.now(), load_id, document_id),
            )
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            cursor.close()
            if not self.db_connection:
                conn.close()

    def _mark_timeout(self, attachment_id, load_id):
        """Mark document extraction as timed out.

        Note: attachment_id here is actually document_id (mapped in _fetch_pending).
        """
        document_id = attachment_id
        conn = self.db_connection or get_connection()
        cursor = conn.cursor()
        try:
            cursor.execute(
                "UPDATE attachment_document SET "
                "extraction_status = 'timeout', extraction_retry_count = extraction_retry_count + 1, "
                "extracted_at = %s, last_load_id = %s "
                "WHERE document_id = %s",
                (datetime.now(), load_id, document_id),
            )
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            cursor.close()
            if not self.db_connection:
                conn.close()

    # ------------------------------------------------------------------
    # Delegating wrappers for backward compatibility
    # ------------------------------------------------------------------
    # These thin wrappers allow tests and callers that reference
    # extractor._extract_pdf(path) etc. to keep working even though
    # the real implementations are now module-level functions.

    @staticmethod
    def _extract_pdf(file_path):
        return _extract_pdf(file_path)

    @staticmethod
    def _extract_pdf_pages(file_path, page_range, table_timeout=10):
        return _extract_pdf_pages(file_path, page_range, table_timeout)

    @staticmethod
    def _extract_docx(file_path):
        return _extract_docx(file_path)

    @staticmethod
    def _extract_xlsx(file_path):
        return _extract_xlsx(file_path)

    @staticmethod
    def _extract_xls(file_path):
        return _extract_xls(file_path)

    @staticmethod
    def _extract_doc(file_path):
        return _extract_doc(file_path)

    @staticmethod
    def _extract_pptx(file_path):
        return _extract_pptx(file_path)

    @staticmethod
    def _extract_rtf(file_path):
        return _extract_rtf(file_path)

    @staticmethod
    def _extract_odt(file_path):
        return _extract_odt(file_path)

    @staticmethod
    def _extract_plain_text(file_path):
        return _extract_plain_text(file_path)

    @staticmethod
    def _resolve_handler(content_type, ext):
        return _resolve_handler(content_type, ext)

    @staticmethod
    def _detect_type_by_magic(file_path):
        return _detect_type_by_magic(file_path)

    @staticmethod
    def _check_ole2_encryption(file_path):
        return _check_ole2_encryption(file_path)

    @staticmethod
    def _process_pdf_blocks(blocks, table_rects=None):
        return _process_pdf_blocks(blocks, table_rects)

    @staticmethod
    def _rect_overlaps_any(x0, y0, x1, y1, rects):
        return _rect_overlaps_any(x0, y0, x1, y1, rects)

    @staticmethod
    def _pdf_table_to_markdown(table):
        return _pdf_table_to_markdown(table)

    @staticmethod
    def _docx_paragraph_to_md(para):
        return _docx_paragraph_to_md(para)

    @staticmethod
    def _docx_table_to_md(table):
        return _docx_table_to_md(table)

    @staticmethod
    def _pptx_table_to_md(table):
        return _pptx_table_to_md(table)

    @staticmethod
    def _odt_table_to_md(table_element):
        return _odt_table_to_md(table_element)

    @staticmethod
    def _find_tables_with_timeout(file_path, page_idx, timeout=10):
        return _find_tables_with_timeout(file_path, page_idx, timeout)

    @staticmethod
    def _cell_str(value):
        return _cell_str(value)


# ------------------------------------------------------------------
# Internal exception types (not exported)
# ------------------------------------------------------------------

class _UnsupportedType(Exception):
    """Raised when the attachment content type is not supported."""


class _SkippedAttachment(Exception):
    """Raised when an attachment is skipped (e.g., already extracted)."""
